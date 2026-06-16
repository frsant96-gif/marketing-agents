from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Cores Solveplan
AZUL = RGBColor(0x00, 0x2F, 0x6C)      # Azul escuro
AMARELO = RGBColor(0xF5, 0xA8, 0x00)   # Amarelo/dourado
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF4, 0xF4, 0xF4)
CINZA_TEXTO = RGBColor(0x44, 0x44, 0x44)
AZUL_CLARO = RGBColor(0xE8, 0xF0, 0xFE)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=14, bold=False, color=RGBColor(0,0,0), align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_divider(slide, left, top, width, color=AMARELO, height=0.05):
    add_rect(slide, left, top, width, height, color)

def slide_layout(prs):
    return prs.slide_layouts[6]  # blank

# =======================
# SLIDE 1 — CAPA
# =======================
slide = prs.slides.add_slide(slide_layout(prs))

# Fundo azul
add_rect(slide, 0, 0, 13.33, 7.5, AZUL)
# Barra amarela inferior
add_rect(slide, 0, 6.8, 13.33, 0.7, AMARELO)
# Barra amarela topo
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)

add_text(slide, "Planejamento de Evento", 0.8, 1.5, 8, 1.2, 36, bold=True, color=BRANCO)
add_text(slide, "SAP BDC Innovation Day 2026", 0.8, 2.7, 10, 0.8, 24, bold=False, color=AMARELO)
add_text(slide, "Dados como Ativo Estratégico", 0.8, 3.3, 10, 0.6, 18, color=BRANCO)
add_text(slide, "26 de agosto de 2026  |  SAP Brasil, São Paulo  |  08:30 às 12h", 0.8, 4.1, 11, 0.5, 14, color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide, "solveplan", 0.8, 6.85, 3, 0.5, 14, bold=True, color=AZUL)

# =======================
# SLIDE 2 — DETALHES E OBJETIVOS
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Detalhes", 0.5, 0.3, 5, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 5.5)

detalhes = [
    ("Nome:", "SAP BDC Innovation Day 2026"),
    ("Tema:", "SAP Business Data Cloud (BDC)"),
    ("Data:", "26 de agosto de 2026"),
    ("Horário:", "08:30 às 12h"),
    ("Local:", "SAP Brasil — Av. das Nações Unidas, 14171"),
    ("Auditório:", "MASP — até 60 pessoas"),
    ("Público-alvo convidado:", "120 pessoas (considerando quebra ~42%)"),
]
y = 1.1
for label, valor in detalhes:
    add_text(slide, label, 0.5, y, 2.5, 0.32, 11, bold=True, color=AZUL)
    add_text(slide, valor, 2.9, y, 5.5, 0.32, 11, color=CINZA_TEXTO)
    y += 0.33

add_text(slide, "Objetivos", 7.5, 0.3, 5, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 7.5, 0.95, 5.3)

objetivos = [
    "Awareness e posicionamento como referência BDC",
    "Expandir contatos e networking estratégico",
    "Gerar novas oportunidades de negócio",
    "Retenção e upsell de clientes ativos",
    "Geração de deals no pipeline",
]
y = 1.1
for obj in objetivos:
    add_text(slide, f"• {obj}", 7.5, y, 5.5, 0.35, 12, color=CINZA_TEXTO)
    y += 0.4

# =======================
# SLIDE 3 — OVERVIEW
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Overview", 0.5, 0.3, 5, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 12.3)

cards = [
    ("Solução", "SAP Business\nData Cloud (BDC)"),
    ("Objetivo", "Awareness e geração\nde demanda (deals no pipe)"),
    ("Etapa do Funil", "Meio e fundo"),
    ("Período", "Jun – Ago 2026"),
    ("Target", "+R$ 2B faturamento"),
    ("Indústria", "Cross industry"),
    ("Persona", "Controller, Head de Dados,\nCIO, CFO — Diretor, Gerente,\nCoordenador, Analista Sênior"),
    ("Conversão", "Reunião / Diagnóstico / PoC"),
]

positions = [
    (0.4, 1.2), (3.5, 1.2), (6.6, 1.2), (9.7, 1.2),
    (0.4, 3.8), (3.5, 3.8), (6.6, 3.8), (9.7, 3.8),
]
for i, (titulo, conteudo) in enumerate(cards):
    x, y = positions[i]
    add_rect(slide, x, y, 2.8, 2.2, AZUL_CLARO)
    add_text(slide, titulo, x+0.15, y+0.1, 2.5, 0.4, 11, bold=True, color=AZUL)
    add_divider(slide, x+0.15, y+0.5, 2.5, AMARELO, 0.04)
    add_text(slide, conteudo, x+0.15, y+0.6, 2.5, 1.4, 10, color=CINZA_TEXTO)

# =======================
# SLIDE 4 — METAS DE SUCESSO
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Metas de Sucesso", 0.5, 0.3, 8, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 12.3)
add_text(slide, "Baseline: edição de fev/2025 — 108 inscrições / 63 presentes / 5 deals / ROI 9.824% / R$ 5.034,80 investimento", 0.5, 1.0, 12.3, 0.4, 10, italic=True, color=CINZA_TEXTO)

metas = [
    ("120", "Inscrições"),
    ("70", "Participantes"),
    ("≥58%", "Taxa de Comparecimento"),
    ("70", "MQLs"),
    ("7", "SQLs"),
    ("6", "Oportunidades"),
    ("R$ 3M", "Pipeline Esperado"),
    (">9.000%", "ROI Estimado"),
]
positions_m = [
    (0.4, 1.5), (2.95, 1.5), (5.5, 1.5), (8.05, 1.5), (10.6, 1.5),
    (0.4, 4.0), (3.8, 4.0), (7.2, 4.0),
]
widths = [2.3, 2.3, 2.3, 2.3, 2.3, 3.0, 3.0, 3.0]
for i, (numero, label) in enumerate(metas):
    x, y = positions_m[i]
    w = widths[i]
    add_rect(slide, x, y, w, 2.0, AZUL)
    add_text(slide, numero, x, y+0.3, w, 0.9, 28, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y+1.15, w, 0.6, 11, color=BRANCO, align=PP_ALIGN.CENTER)

# =======================
# SLIDE 5 — AGENDA
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Agenda", 0.5, 0.3, 5, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 12.3)

agenda = [
    ("08:30 – 09:00", "Coffee de boas-vindas + networking", AMARELO),
    ("09:00 – 09:30", "Abertura — Solveplan + contexto SAP BDC no mercado", AZUL),
    ("09:30 – 10:15", "SAP Business Data Cloud: unificando dados SAP e non-SAP com IA integrada + demo ao vivo", AZUL),
    ("10:15 – 10:45", "Case de sucesso — cliente Solveplan (a confirmar: VALE / Klabin / Aegea)", AZUL),
    ("10:45 – 11:15", "Mesa redonda: desafios reais de dados nas empresas", AZUL),
    ("11:15 – 11:30", "Encerramento + próximos passos", AZUL),
    ("11:30 – 12:00", "Networking livre", AMARELO),
]

y = 1.1
for horario, atividade, cor in agenda:
    add_rect(slide, 0.5, y, 2.5, 0.62, cor)
    add_text(slide, horario, 0.5, y, 2.5, 0.62, 11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.1, y, 9.8, 0.62, AZUL_CLARO if cor == AZUL else RGBColor(0xFF, 0xF0, 0xCC))
    add_text(slide, atividade, 3.25, y+0.05, 9.5, 0.52, 11, color=CINZA_TEXTO)
    y += 0.72

# =======================
# SLIDE 6 — ORÇAMENTO
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Orçamento", 0.5, 0.3, 5, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 12.3)

# Tabela orçamento
headers = ["Item", "Valor Previsto", "Realizado"]
col_widths = [6.0, 3.0, 3.0]
col_starts = [0.5, 6.6, 9.7]
y = 1.1
add_rect(slide, 0.5, y, 12.3, 0.45, AMARELO)
for i, h in enumerate(headers):
    add_text(slide, h, col_starts[i]+0.1, y+0.05, col_widths[i]-0.2, 0.35, 12, bold=True, color=AZUL)
y += 0.45

itens = [
    ("Coffee (70 pax)", "R$ 4.400,00", ""),
    ("Brindes (Moleskine + Caneta)", "R$ 800,00", ""),
    ("LinkedIn Ads (convite para o evento)", "R$ 800,00", ""),
    ("Total", "R$ 6.000,00", ""),
]
for i, (item, previsto, real) in enumerate(itens):
    cor_bg = CINZA_CLARO if i % 2 == 0 else BRANCO
    is_total = item == "Total"
    if is_total:
        cor_bg = AZUL
    add_rect(slide, 0.5, y, 12.3, 0.45, cor_bg)
    text_color = BRANCO if is_total else CINZA_TEXTO
    add_text(slide, item, col_starts[0]+0.1, y+0.05, col_widths[0]-0.2, 0.35, 11, bold=is_total, color=text_color)
    add_text(slide, previsto, col_starts[1]+0.1, y+0.05, col_widths[1]-0.2, 0.35, 11, bold=is_total, color=text_color)
    add_text(slide, real, col_starts[2]+0.1, y+0.05, col_widths[2]-0.2, 0.35, 11, bold=is_total, color=text_color)
    y += 0.45

y += 0.3
add_text(slide, "Custo por inscrito esperado:", 0.5, y, 4, 0.35, 11, bold=True, color=AZUL)
add_text(slide, "~R$ 50,00", 4.6, y, 3, 0.35, 11, color=CINZA_TEXTO)
y += 0.4
add_text(slide, "Custo por participante esperado:", 0.5, y, 4, 0.35, 11, bold=True, color=AZUL)
add_text(slide, "~R$ 86,00", 4.6, y, 3, 0.35, 11, color=CINZA_TEXTO)
y += 0.4
add_text(slide, "Custo por oportunidade esperado:", 0.5, y, 4, 0.35, 11, bold=True, color=AZUL)
add_text(slide, "~R$ 1.000,00", 4.6, y, 3, 0.35, 11, color=CINZA_TEXTO)

# Nota comparativa
add_rect(slide, 0.5, 5.5, 12.3, 0.8, AZUL_CLARO)
add_text(slide, "Referência fev/2025: R$ 5.034,80 investidos → 5 deals → ROI 9.824%. Com R$ 6.000,00 e meta de 6 deals, projeção de ROI superior.",
         0.65, 5.55, 12.0, 0.7, 10, italic=True, color=AZUL)

# =======================
# SLIDE 7 — CRONOGRAMA
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Cronograma", 0.5, 0.3, 5, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 12.3)

fases = [
    ("JUN 2026", "Planejamento\nDefinições e confirmações"),
    ("JUL 2026", "Criação e Divulgação\nMateriais + início dos convites (21/07)"),
    ("AGO 2026", "Convites e RSVP\nReforço WhatsApp + confirmações"),
    ("26/08", "EVENTO"),
    ("SET 2026", "Follow-up e Report\nCRM + relatório de resultados"),
]

x = 0.5
for i, (mes, desc) in enumerate(fases):
    is_evento = mes == "26/08"
    cor = AMARELO if is_evento else AZUL
    add_rect(slide, x, 1.2, 2.3, 0.5, cor)
    add_text(slide, mes, x, 1.2, 2.3, 0.5, 12, bold=True, color=AZUL if is_evento else BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, x+1.15, 1.7, 0.05, 3.5, RGBColor(0xCC, 0xCC, 0xCC))
    add_rect(slide, x, 1.7, 2.3, 3.2, AZUL_CLARO if not is_evento else RGBColor(0xFF, 0xF0, 0xCC))
    add_text(slide, desc, x+0.1, 1.8, 2.1, 2.8, 10, color=CINZA_TEXTO if not is_evento else AZUL)
    x += 2.56

# Detalhe emails
add_text(slide, "Sequência de e-mails: 21/07 → 28/07 → 04/08 → 11/08 → 18/08 → Lembrete 24/08 → Agradecimento 27/08",
         0.5, 5.1, 12.3, 0.4, 10, italic=True, color=CINZA_TEXTO)

# =======================
# SLIDE 8 — COMUNICAÇÃO
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Comunicação", 0.5, 0.3, 5, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 12.3)

# Coluna esquerda — e-mails
add_text(slide, "E-mails", 0.5, 1.1, 5.5, 0.4, 14, bold=True, color=AZUL)
emails = [
    ("E-mail 01", "21/07", "1º convite"),
    ("E-mail 02", "28/07", "2º convite"),
    ("E-mail 03", "04/08", "3º convite"),
    ("E-mail 04", "11/08", "4º convite + WhatsApp"),
    ("E-mail 05", "18/08", "RSVP / últimas vagas"),
    ("Lembrete",  "24/08", "D-2"),
    ("Agradecimento", "27/08", "Pós-evento"),
]
y = 1.55
for i, (email, data, obj) in enumerate(emails):
    cor_bg = AZUL if i % 2 == 0 else AZUL_CLARO
    cor_txt = BRANCO if i % 2 == 0 else AZUL
    add_rect(slide, 0.5, y, 5.8, 0.42, cor_bg)
    add_text(slide, email, 0.6, y+0.04, 1.8, 0.34, 10, bold=True, color=cor_txt)
    add_text(slide, data, 2.4, y+0.04, 1.0, 0.34, 10, color=cor_txt, align=PP_ALIGN.CENTER)
    add_text(slide, obj, 3.5, y+0.04, 2.7, 0.34, 10, color=cor_txt)
    y += 0.44

# Coluna direita — canais
add_text(slide, "Canais de Divulgação", 7.1, 1.1, 5.8, 0.4, 14, bold=True, color=AZUL)
canais = [
    ("LinkedIn", "Posts orgânicos + anúncios pagos (R$ 800,00)"),
    ("Instagram", "Posts + Stories + Reels durante o evento"),
    ("WhatsApp", "Disparo direto para contatos qualificados"),
    ("E-mail Marketing", "Base de contatos segmentada"),
    ("Site SAP Brasil", "Divulgação na plataforma da SAP"),
]
y = 1.55
for canal, desc in canais:
    add_rect(slide, 7.1, y, 5.8, 0.75, AZUL_CLARO)
    add_text(slide, canal, 7.25, y+0.05, 5.5, 0.3, 11, bold=True, color=AZUL)
    add_text(slide, desc, 7.25, y+0.35, 5.5, 0.32, 10, color=CINZA_TEXTO)
    y += 0.85

# =======================
# SLIDE 9 — ENGAJAMENTO E BRINDES
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Engajamento, Brindes e Vestimenta", 0.5, 0.3, 10, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 12.3)

# Pesquisa
add_rect(slide, 0.5, 1.1, 5.8, 2.5, AZUL_CLARO)
add_text(slide, "Pesquisa & Engajamento", 0.65, 1.2, 5.5, 0.4, 13, bold=True, color=AZUL)
add_divider(slide, 0.65, 1.62, 5.5, AMARELO, 0.04)
pesq = [
    "QR Code de Pesquisa de Satisfação",
    "Plataforma: Qualtrics (digital)",
    "Brinde condicionado ao preenchimento via QR Code",
    "Captação de dados para CRM no mesmo momento",
]
y = 1.7
for p in pesq:
    add_text(slide, f"• {p}", 0.65, y, 5.5, 0.35, 11, color=CINZA_TEXTO)
    y += 0.38

# Brindes
add_rect(slide, 0.5, 3.8, 5.8, 2.0, AZUL_CLARO)
add_text(slide, "Brindes", 0.65, 3.9, 5.5, 0.4, 13, bold=True, color=AZUL)
add_divider(slide, 0.65, 4.32, 5.5, AMARELO, 0.04)
add_text(slide, "• Moleskine (caderneta)", 0.65, 4.4, 5.5, 0.35, 11, color=CINZA_TEXTO)
add_text(slide, "• Caneta", 0.65, 4.78, 5.5, 0.35, 11, color=CINZA_TEXTO)
add_text(slide, "Entrega mediante confirmação de preenchimento da pesquisa", 0.65, 5.15, 5.5, 0.45, 10, italic=True, color=CINZA_TEXTO)

# Vestimenta
add_rect(slide, 7.0, 1.1, 5.8, 2.0, AZUL_CLARO)
add_text(slide, "Vestimenta do Time", 7.15, 1.2, 5.5, 0.4, 13, bold=True, color=AZUL)
add_divider(slide, 7.15, 1.62, 5.5, AMARELO, 0.04)
add_text(slide, "Presença no evento:", 7.15, 1.7, 5.5, 0.35, 11, bold=True, color=AZUL)
add_text(slide, "Sócios + Francielle Beline", 7.15, 2.05, 5.5, 0.35, 11, color=CINZA_TEXTO)
add_text(slide, "Vestimenta:", 7.15, 2.45, 5.5, 0.35, 11, bold=True, color=AZUL)
add_text(slide, "Camiseta Insider com layout Solveplan", 7.15, 2.8, 5.5, 0.35, 11, color=CINZA_TEXTO)

# Coffee
add_rect(slide, 7.0, 3.3, 5.8, 2.5, AZUL_CLARO)
add_text(slide, "Coffee Break", 7.15, 3.4, 5.5, 0.4, 13, bold=True, color=AZUL)
add_divider(slide, 7.15, 3.82, 5.5, AMARELO, 0.04)
coffee_items = [
    "Bebidas: café, leite, suco, chá, água, iogurte",
    "Doces + frutas: mini bolo, waffle, salada de frutas",
    "Salgados: mini sanduíches, pão de queijo, folhado",
    "70 pessoas | ~R$ 4.400,00 | 2 funcionários",
]
y = 3.9
for c in coffee_items:
    add_text(slide, f"• {c}", 7.15, y, 5.5, 0.35, 10, color=CINZA_TEXTO)
    y += 0.38

# =======================
# SLIDE 10 — RESULTADOS ESPERADOS E REFERÊNCIA
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.12, AZUL)

add_text(slide, "Resultados Esperados", 0.5, 0.3, 10, 0.6, 22, bold=True, color=AZUL)
add_divider(slide, 0.5, 0.95, 12.3)

# Cards comparativos
comparativo = [
    ("Inscrições", "108", "120"),
    ("Participantes", "63", "70"),
    ("Taxa comp.", "58%", "≥ 58%"),
    ("MQLs", "60", "70"),
    ("SQLs", "6", "7"),
    ("Deals", "5", "6"),
    ("Investimento", "R$5.034", "R$6.000"),
    ("ROI", "9.824%", ">9.000%"),
]

add_text(slide, "Fev/2025", 3.5, 1.05, 3.5, 0.35, 11, bold=True, color=CINZA_TEXTO, align=PP_ALIGN.CENTER)
add_text(slide, "Meta Ago/2026", 7.5, 1.05, 4.0, 0.35, 11, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

y = 1.45
for i, (metrica, anterior, meta) in enumerate(comparativo):
    cor_bg = CINZA_CLARO if i % 2 == 0 else BRANCO
    add_rect(slide, 0.5, y, 12.3, 0.52, cor_bg)
    add_text(slide, metrica, 0.6, y+0.08, 2.8, 0.36, 11, bold=True, color=AZUL)
    add_text(slide, anterior, 3.5, y+0.08, 3.5, 0.36, 11, color=CINZA_TEXTO, align=PP_ALIGN.CENTER)
    add_rect(slide, 7.4, y+0.06, 4.0, 0.40, AZUL)
    add_text(slide, meta, 7.4, y+0.08, 4.0, 0.36, 11, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
    y += 0.54

# =======================
# SLIDE 11 — PRÓXIMOS PASSOS
# =======================
slide = prs.slides.add_slide(slide_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, AZUL)
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_rect(slide, 0, 6.8, 13.33, 0.7, AMARELO)

add_text(slide, "Próximos Passos", 0.8, 0.4, 10, 0.7, 28, bold=True, color=BRANCO)
add_divider(slide, 0.8, 1.15, 11.7, AMARELO)

passos = [
    ("1", "Confirmar reserva do Auditório MASP na SAP Brasil", "até 20/06"),
    ("2", "Definir nome definitivo do evento", "até 20/06"),
    ("3", "Confirmar case de sucesso com cliente (VALE / Klabin / Aegea)", "até 30/06"),
    ("4", "Produção dos materiais (convite, e-mails, anúncios, card WhatsApp)", "01/07 – 14/07"),
    ("5", "Início dos convites — E-mail 01 + LinkedIn Ads", "21/07"),
    ("6", "Disparo WhatsApp para base qualificada", "30/07 – 04/08"),
    ("7", "Execução do evento", "26/08"),
    ("8", "Follow-up e relatório de resultados", "27/08 – 05/09"),
]

y = 1.35
for num, acao, prazo in passos:
    add_rect(slide, 0.8, y, 0.55, 0.5, AMARELO)
    add_text(slide, num, 0.8, y, 0.55, 0.5, 14, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    add_text(slide, acao, 1.5, y+0.06, 9.0, 0.38, 12, color=BRANCO)
    add_text(slide, prazo, 10.6, y+0.06, 2.4, 0.38, 11, color=AMARELO, align=PP_ALIGN.RIGHT)
    y += 0.6

add_text(slide, "solveplan", 0.8, 6.85, 4, 0.45, 13, bold=True, color=AZUL)

# Salvar
output_path = r"c:\Users\franc\solveplan.com\Roberto Molina - Marketing\1. MKT Estrategy\3. Agentes de IA\ccos-ratos\eventos\sap-bdc-innovation-day-2026\SAP_BDC_Innovation_Day_2026_Planejamento.pptx"
prs.save(output_path)
print(f"PPT salvo: {output_path}")