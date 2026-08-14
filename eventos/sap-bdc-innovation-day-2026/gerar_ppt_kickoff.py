from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Cores Solveplan
AZUL = RGBColor(0x00, 0x2F, 0x6C)
AMARELO = RGBColor(0xF5, 0xA8, 0x00)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF4, 0xF4, 0xF4)
CINZA_TEXTO = RGBColor(0x44, 0x44, 0x44)
VERMELHO = RGBColor(0xC0, 0x2B, 0x2B)
VERDE = RGBColor(0x1E, 0x7B, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def slide_layout():
    return prs.slide_layouts[6]

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=14, bold=False, color=RGBColor(0,0,0), align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_bullets(slide, items, left, top, width, height, font_size=13, color=CINZA_TEXTO, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox

def header(slide, titulo, kicker=None):
    add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
    add_rect(slide, 0, 7.38, 13.33, 0.12, AZUL)
    if kicker:
        add_text(slide, kicker, 0.5, 0.28, 8, 0.4, 13, bold=True, color=AMARELO)
        add_text(slide, titulo, 0.5, 0.6, 10, 0.6, 24, bold=True, color=AZUL)
    else:
        add_text(slide, titulo, 0.5, 0.35, 10, 0.6, 24, bold=True, color=AZUL)
    add_rect(slide, 0.5, 1.15, 2.2, 0.04, AMARELO)

def add_table(slide, headers, rows, left, top, col_widths, row_height=0.4, font_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_width = sum(col_widths)
    total_height = row_height * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(total_width), Inches(total_height))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = BRANCO
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CINZA_CLARO if r % 2 == 0 else BRANCO
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = CINZA_TEXTO
    return table_shape

# =======================
# SLIDE 1 — CAPA
# =======================
slide = prs.slides.add_slide(slide_layout())
add_rect(slide, 0, 0, 13.33, 7.5, AZUL)
add_rect(slide, 0, 6.8, 13.33, 0.7, AMARELO)
add_rect(slide, 0, 0, 13.33, 0.12, AMARELO)
add_text(slide, "Kick Off", 0.8, 1.5, 8, 1.2, 40, bold=True, color=BRANCO)
add_text(slide, "SAP Innovation Day | Data Foundation First", 0.8, 2.6, 11, 0.7, 22, color=AMARELO)
add_text(slide, "Antes da IA, vêm os dados", 0.8, 3.2, 10, 0.6, 16, italic=True, color=BRANCO)
add_text(slide, "Reunião de kick off — 14/08 · 16h  |  Evento em 26/08 · SAP Brasil", 0.8, 4.0, 11, 0.5, 14, color=RGBColor(0xCC,0xCC,0xCC))
add_text(slide, "solveplan", 0.8, 6.85, 3, 0.5, 14, bold=True, color=AZUL)

# =======================
# SLIDE 2 — VISÃO GERAL DO EVENTO
# =======================
slide = prs.slides.add_slide(slide_layout())
header(slide, "Visão Geral do Evento", "1. ABERTURA")
info = [
    ("Data", "26/08 · 08h30 às 12h"),
    ("Local", "SAP Brasil — Av. das Nações Unidas, 14171"),
    ("Tema", "SAP Business Data Cloud com IA"),
    ("Case de sucesso", "✅ SSA Alimentos — apresentado por Rodrigo Dumont"),
    ("Público-alvo", "50/60 pessoas convidadas (considerando quebra)"),
    ("Objetivo", "Gerar demanda qualificada para SAP BDC, aproximar contas\nestratégicas (+R$800M/ano) e converter em reunião/diagnóstico/PoC"),
]
top = 1.5
for label, val in info:
    add_text(slide, label, 0.5, top, 2, 0.4, 14, bold=True, color=AZUL)
    add_text(slide, val, 2.7, top, 9.5, 0.7, 14, color=CINZA_TEXTO)
    top += 0.9

# =======================
# SLIDE 3 — METAS
# =======================
slide = prs.slides.add_slide(slide_layout())
header(slide, "Metas do Evento", "2. METAS")
metas = [
    ("120", "Inscrições"),
    ("50", "Participantes"),
    ("≥42%", "Comparecimento"),
    ("20", "Contas"),
    ("5", "Oportunidades"),
    ("R$ 1,1M", "Pipeline potencial"),
    ("R$ 6,5 mil", "Investimento"),
    ("168x", "ROI potencial"),
]
cols = 4
box_w, box_h, gap_x, gap_y = 2.8, 1.5, 0.25, 0.3
start_x, start_y = 0.6, 1.6
for i, (num, label) in enumerate(metas):
    r, c = divmod(i, cols)
    x = start_x + c * (box_w + gap_x)
    y = start_y + r * (box_h + gap_y)
    color = AZUL if i < 6 else AMARELO
    add_rect(slide, x, y, box_w, box_h, color)
    add_text(slide, num, x, y+0.2, box_w, 0.6, 26, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y+0.85, box_w, 0.5, 13, color=BRANCO, align=PP_ALIGN.CENTER)

# =======================
# SLIDE 3B — O QUE TEREMOS
# =======================
slide = prs.slides.add_slide(slide_layout())
header(slide, "O Que Teremos", "CONFIRMADO")
add_bullets(slide, [
    "Coffee",
    "Brindes",
    "Apresentação de case (SSA Alimentos)",
    "Pesquisa de engajamento",
], 0.6, 1.6, 10, 3, 16, line_spacing=1.4)
add_text(slide, "*Inscrições até o momento: XX — atualizar na reunião", 0.6, 5.5, 10, 0.5, 13, italic=True, color=RGBColor(0x99,0x66,0x00))

# =======================
# SLIDE 4 — ALERTA DE PRAZO
# =======================
slide = prs.slides.add_slide(slide_layout())
header(slide, "Onde Estamos", "3. STATUS")
add_rect(slide, 0.5, 1.4, 12.3, 1.0, RGBColor(0xFC, 0xEA, 0xD8))
add_text(slide, "⚠ ALERTA DE PRAZO", 0.8, 1.55, 5, 0.4, 14, bold=True, color=VERMELHO)
add_text(slide, "Hoje é 14/08 — evento em 12 dias. Case e palestrantes já confirmados.\nDivulgação e e-mails de convite seguem atrasados — foco #1 da reunião.",
         0.8, 1.9, 11.5, 0.6, 13, color=CINZA_TEXTO)

headers_tbl = ["Frente", "Status", "Pendência crítica"]
rows = [
    ["Case de sucesso", "✅ SSA Alimentos / Rodrigo Dumont", "Alinhar roteiro do case"],
    ["Palestrantes", "✅ Junior Freitas, Andrey, André Ferreira", "Alinhar conteúdo e tempo de fala"],
    ["Lista de convidados", "Ajustada p/ 50/60 pessoas", "Fechar lista final hoje/amanhã"],
    ["Inscrições até agora", "Em aberto (XX)", "Levantar número atual"],
    ["Página + formulário", "Layout pronto", "Confirmar publicação/link ativo"],
    ["E-mails de convite", "Cadência atrasada", "Replanejar disparos comprimidos"],
    ["Posts redes sociais", "Atrasados (post 01 e 02)", "Publicar retroativo esta semana"],
    ["Anúncios (LinkedIn/Google)", "Budget aprovado R$1.000", "Confirmar campanha no ar"],
    ["Brindes", "Pendente aprovação", "Aprovar R$ 300,00"],
]
add_table(slide, headers_tbl, rows, 0.5, 2.6, [3.2, 3.5, 5.4], row_height=0.5, font_size=11)

# =======================
# SLIDE 5 — AGENDA DO DIA
# =======================
slide = prs.slides.add_slide(slide_layout())
header(slide, "Agenda do Dia 26/08", "4. PROGRAMAÇÃO")
headers_ag = ["Horário", "Bloco", "Responsável"]
rows_ag = [
    ["08h30–09h00", "Recepção e networking executivo", "Coffee / chegada"],
    ["09h00–09h30", "Abertura: antes da IA, vêm os dados", "Junior Freitas e Andrey"],
    ["09h30–10h00", "SAP BDC: fundação de dados integrados", "SAP, Andrey e André Ferreira"],
    ["10h00–10h30", "Demonstração prática", "Andrey e André Ferreira"],
    ["10min", "Break", "All"],
    ["10h40–11h10", "Case de sucesso SSA", "Rodrigo Dumont (SSA Alimentos)"],
    ["11h10–11h40", "Próximos passos / encerramento", "—"],
    ["11h40–12h00", "Networking final", "—"],
]
add_table(slide, headers_ag, rows_ag, 0.5, 1.5, [2.5, 6.0, 3.6], row_height=0.55, font_size=12)

# =======================
# SLIDE 6 — MENSAGEM-CHAVE
# =======================
slide = prs.slides.add_slide(slide_layout())
header(slide, "Mensagem-Chave", "5. POSICIONAMENTO")
add_rect(slide, 0.5, 1.5, 12.3, 1.1, AZUL)
add_text(slide, '"Sua empresa já tem os dados. O SAP BDC entrega a IA."', 0.8, 1.75, 11.5, 0.6, 20, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_bullets(slide, [
    "Tom: expert/conselheiro, direto, orientado a negócio",
    "Não mencionar preço sem qualificação",
    "Não comparar com Databricks/Snowflake sem contexto adequado",
    "CTA lead: inscrição no evento",
    "CTA SQL: diagnóstico gratuito / reunião com especialista",
], 0.8, 3.0, 11, 3.0, 15)

# =======================
# SLIDE 7 — DECISÕES DE HOJE
# =======================
slide = prs.slides.add_slide(slide_layout())
header(slide, "Decisões para Sair Desta Reunião", "6. AÇÃO")
decisoes = [
    "Reportar número de inscrições até hoje (comparar com meta de 120)",
    "Fechar lista final de 50/60 convidados",
    "Alinhar roteiro do case com Rodrigo Dumont (SSA Alimentos)",
    "Replanejar cadência de e-mails (disparos comprimidos até 24/08)",
    "Quem publica os posts atrasados (lançamento + convite 01) esta semana",
    "Confirmar se anúncios (LinkedIn/Google) já estão ativos",
    "Aprovar custo de brindes (R$ 300)",
    "Alinhar conteúdo e tempo de fala de Junior Freitas, Andrey e André Ferreira",
]
add_bullets(slide, decisoes, 0.6, 1.5, 12, 5, 15, line_spacing=1.3)

# =======================
# SLIDE 8 — PRÓXIMOS PASSOS
# =======================
slide = prs.slides.add_slide(slide_layout())
header(slide, "Próximos Passos Imediatos", "7. CHECKLIST")
passos = [
    "Reportar número de inscrições atuais",
    "Fechar lista final de convidados (50/60)",
    "Alinhar roteiro do case com Rodrigo Dumont",
    "Disparar e-mail de convite (comprimido) ainda esta semana",
    "Publicar post de lançamento + convite 01 (atrasados)",
    "Validar página do evento e formulário ao vivo",
    "Confirmar campanha de anúncios no ar",
    "Aprovar budget de brindes",
    "Alinhar apresentação com os 3 palestrantes internos",
]
add_bullets(slide, passos, 0.6, 1.5, 12, 5, 15, line_spacing=1.3)
add_text(slide, "Fonte: Planejamento de Evento — Abril 2026 | eventos/sap-bdc-innovation-day-2026/", 0.6, 6.9, 11, 0.4, 10, italic=True, color=RGBColor(0x99,0x99,0x99))

prs.save("SAP_BDC_Kickoff_16h_2026-08-14.pptx")
print("OK")
