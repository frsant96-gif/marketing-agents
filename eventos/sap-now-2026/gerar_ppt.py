from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Cores do template Solveplan ──────────────────────────────
BG_LIGHT   = RGBColor(0xF0, 0xF0, 0xF2)   # fundo claro (cinza off-white)
BG_DARK    = RGBColor(0x0A, 0x08, 0x37)   # fundo escuro (azul marinho)
AZUL_VIVO  = RGBColor(0x00, 0x6A, 0xFF)   # azul destaque / títulos
NAVY       = RGBColor(0x0A, 0x08, 0x37)   # azul marinho (texto sobre fundo claro)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_TXT  = RGBColor(0x44, 0x44, 0x55)   # texto secundário
VERDE_HDR  = RGBColor(0x7F, 0xFF, 0x96)   # cabeçalho tabela verde (como template)
AZUL_HDR   = RGBColor(0x00, 0x6A, 0xFF)   # cabeçalho tabela azul
AZUL_ROW1  = RGBColor(0xD6, 0xE8, 0xFF)   # linha par — azul claro
AZUL_ROW2  = RGBColor(0xEB, 0xF3, 0xFF)   # linha ímpar — azul mais claro
AMARELO    = RGBColor(0xFF, 0xC0, 0x00)   # atenção / alerta

W = Inches(13.33)
H = Inches(7.5)

LOGO_DARK  = r"c:\Users\franc\solveplan.com\Roberto Molina - Marketing\1. MKT Estrategy\3. Agentes de IA\ccos-ratos\marca\logo-escuro.png.png"
LOGO_LIGHT = r"c:\Users\franc\solveplan.com\Roberto Molina - Marketing\1. MKT Estrategy\3. Agentes de IA\ccos-ratos\marca\logo-claro.png.png"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

slide_num = [0]


# ── Helpers ──────────────────────────────────────────────────

def new_slide(dark=False):
    s = prs.slides.add_slide(blank)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK if dark else BG_LIGHT
    slide_num[0] += 1
    return s


def txb(slide, text, x, y, w, h, size=14, bold=False,
        color=NAVY, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return tb


def box(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def title_block(slide, title, dark=False):
    """Título + linha decorativa abaixo — padrão do template."""
    tc = AZUL_VIVO if dark else AZUL_VIVO
    txb(slide, title, Inches(0.55), Inches(0.25), Inches(12.5), Inches(0.7),
        size=22, bold=True, color=tc)
    # linha decorativa
    ln = box(slide, Inches(0.55), Inches(0.95), Inches(1.1), Inches(0.045),
             NAVY if not dark else BRANCO)


def footer(slide, dark=False, page_label=None):
    """Rodapé: logo + linha + texto + número."""
    logo_file = LOGO_LIGHT if dark else LOGO_DARK
    if os.path.exists(logo_file):
        try:
            slide.shapes.add_picture(logo_file,
                Inches(0.35), Inches(6.98), Inches(0.85), Inches(0.38))
        except Exception:
            pass

    # linha
    ln = box(slide, Inches(1.35), Inches(7.15), Inches(11.6), Inches(0.022),
             AZUL_VIVO if dark else NAVY)

    # texto central
    lbl_color = AZUL_VIVO if dark else AZUL_VIVO
    txb(slide, "Planejamento SAP NOW AI Tour 2026 — Solveplan",
        Inches(1.4), Inches(7.18), Inches(9.5), Inches(0.28),
        size=9, color=lbl_color)

    # número
    num = str(slide_num[0]) if page_label is None else page_label
    txb(slide, num, Inches(12.6), Inches(7.18), Inches(0.6), Inches(0.28),
        size=9, color=lbl_color, align=PP_ALIGN.RIGHT)


def table_block(slide, headers, rows, x, y, col_widths,
                row_h=Inches(0.38), hdr_color=AZUL_HDR):
    """Tabela no estilo do template (cabeçalho colorido + linhas alternadas)."""
    xx = x
    for i, h in enumerate(headers):
        box(slide, xx, y, col_widths[i], Inches(0.44), hdr_color)
        txb(slide, h, xx + Inches(0.06), y + Inches(0.05),
            col_widths[i] - Inches(0.06), Inches(0.38),
            size=11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        xx += col_widths[i]

    for ri, row in enumerate(rows):
        yy = y + Inches(0.44) + ri * row_h
        bg = AZUL_ROW1 if ri % 2 == 0 else AZUL_ROW2
        xx = x
        for ci, cell in enumerate(row):
            box(slide, xx, yy, col_widths[ci], row_h, bg)
            c_bold = cell.startswith("⚠️") or cell.startswith("TOTAL") or cell.startswith("**")
            txb(slide, cell.strip("*"),
                xx + Inches(0.06), yy + Inches(0.03),
                col_widths[ci] - Inches(0.06), row_h - Inches(0.04),
                size=10, bold=c_bold,
                color=NAVY if bg != BG_DARK else BRANCO,
                align=PP_ALIGN.CENTER)
            xx += col_widths[ci]


def kv_block(slide, rows, x, y, w1=Inches(3.2), w2=Inches(9.5)):
    for i, (k, v) in enumerate(rows):
        yy = y + Inches(i * 0.46)
        bg = AZUL_ROW1 if i % 2 == 0 else AZUL_ROW2
        box(slide, x, yy, w1 + w2, Inches(0.42), bg)
        txb(slide, k, x + Inches(0.08), yy + Inches(0.04), w1, Inches(0.38),
            size=11, bold=True, color=NAVY)
        txb(slide, v, x + w1 + Inches(0.08), yy + Inches(0.04), w2, Inches(0.38),
            size=11, color=NAVY)


def bullets(slide, items, x, y, w, size=13, color=NAVY, spacing=0.42):
    for i, item in enumerate(items):
        txb(slide, f"•  {item}", x, y + Inches(i * spacing), w, Inches(0.42),
            size=size, color=color)


def highlight_box(slide, text, x, y, w, h, bg=BG_DARK, tc=BRANCO, size=13):
    box(slide, x, y, w, h, bg)
    txb(slide, text, x + Inches(0.15), y + Inches(0.08),
        w - Inches(0.2), h - Inches(0.1), size=size, color=tc, bold=True)


# ════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════
s = new_slide(dark=True)

# Logo centralizado (grande)
logo_file = LOGO_LIGHT
if os.path.exists(logo_file):
    try:
        s.shapes.add_picture(logo_file,
            Inches(1.5), Inches(2.6), Inches(2.8), Inches(1.25))
    except Exception:
        txb(s, "solveplan", Inches(1.5), Inches(2.8), Inches(3), Inches(0.7),
            size=36, bold=True, color=BRANCO)

# Linha vertical separadora (como no template)
box(s, Inches(4.6), Inches(2.5), Inches(0.025), Inches(1.55), BRANCO)

# Título e info ao lado da linha
txb(s, "SAP NOW AI Tour\nBrazil 2026",
    Inches(4.9), Inches(2.55), Inches(7.5), Inches(1.1),
    size=36, bold=True, color=BRANCO)

txb(s, "Planejamento Solveplan",
    Inches(4.9), Inches(3.75), Inches(7), Inches(0.5),
    size=18, color=AZUL_VIVO)

txb(s, "9 e 10 de setembro de 2026  |  Transamérica Expo Center — São Paulo",
    Inches(4.9), Inches(4.3), Inches(7.5), Inches(0.45),
    size=13, color=RGBColor(0xAA, 0xBB, 0xCC))

footer(s, dark=True, page_label="")


# ════════════════════════════════════════════════════════════
# SLIDE 2 — OBJETIVOS
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Objetivos")

items = [
    "Geração de leads qualificados — meta: 130",
    "Geração de reuniões agendadas — meta: 28",
    "Posicionar Solveplan como referência em SAP Business Data Cloud (BDC)",
    "Fomentar relacionamento com clientes e expandir networking",
    "Gerar pipeline de novos negócios — meta: R$ 3.200.000",
]
bullets(s, items, Inches(0.8), Inches(1.2), Inches(11.5), size=16)

box(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.72), NAVY)
txb(s, '💡  Mensagem-chave: "Seus dados existem — mas você ainda demora dias para fechar o mês. Em 2026, isso é risco. A Solveplan resolve."',
    Inches(0.75), Inches(5.68), Inches(12.0), Inches(0.6),
    size=13, bold=True, color=BRANCO)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — OVERVIEW
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Overview")
kv_block(s, [
    ("Solução",       "SAP Business Data Cloud (BDC) — foco principal"),
    ("Objetivo",      "Geração de demanda + conversão de leads (deals no pipe)"),
    ("Etapa do funil","Meio e fundo"),
    ("Período",       "Jun, Jul, Ago, Set 2026"),
    ("Target",        "+R$ 500M faturamento / cross industry"),
    ("Persona",       "CIO, CFO, Controller, Head de Dados/BI, COO"),
    ("Conversão",     "Evento SAP NOW AI Tour 2026"),
], Inches(0.55), Inches(1.15))
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — DETALHES DO EVENTO
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Detalhes do Evento")
kv_block(s, [
    ("Datas",           "9 e 10 de setembro de 2026"),
    ("Horário",         "Das 8h às 19h"),
    ("Local",           "Transamérica Expo Center — São Paulo"),
    ("Audiência total", "3.900 pessoas"),
    ("Clientes/Prospects", "2.700  |  Decision Makers: 73%"),
    ("Patrocinadores",  "70"),
    ("Cota Solveplan",  "Gold — inclui 1 sessão de conteúdo de 20 min"),
    ("Tema do evento",  "Bring it at SAP NOW AI Tour 2026"),
    ("Nossa posição",   "A confirmar via manual CAEX (recebido em 18/05) — posição não privilegiada"),
], Inches(0.55), Inches(1.15), w1=Inches(3.5), w2=Inches(9.2))
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — METAS E CRITÉRIO DE SUCESSO
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Metas e Critério de Sucesso")
table_block(s,
    ["Métrica", "Meta 2026", "Ref. 2025"],
    [
        ("Leads totais",           "130",           "186"),
        ("Leads qualificados",     "44",            "—"),
        ("Reuniões agendadas",     "28",            "26"),
        ("Oportunidades abertas",  "8",             "—"),
        ("Pipeline gerado",        "R$ 3.200.000",  "—"),
        ("Empresas abordadas SDR", "120+",          "72"),
        ("Agendas via SDR",        "12",            "5"),
        ("Custo por lead",         "R$ 2.127",      "—"),
        ("ROI esperado",           "11,6x",         "—"),
    ],
    Inches(1.5), Inches(1.15),
    [Inches(6.5), Inches(2.8), Inches(2.8)],
    row_h=Inches(0.44),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — STAFF
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Staff")
txb(s, "Time a definir com sócios  |  Ref. 2025: 9 pessoas (8 credenciais patrocinador + 1 staff)",
    Inches(0.55), Inches(1.1), Inches(12.5), Inches(0.35),
    size=12, color=CINZA_TXT, italic=True)

table_block(s,
    ["#", "Nome", "Credencial", "Obs."],
    [
        ("1–8", "A definir", "Patrocinador", ""),
        ("9",   "Fran",      "Staff",        "Organização geral"),
    ],
    Inches(0.55), Inches(1.5),
    [Inches(0.7), Inches(4.5), Inches(3.0), Inches(4.6)],
    row_h=Inches(0.42),
)

txb(s, "Regras de credencial:", Inches(0.55), Inches(2.85), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Até 4 trocas/dia — proibido das 12h às 14h (exige documento)",
    "Retirada dia 08/set no CAEX (8h–18h) ou 09–10/set no credenciamento (8h–17h)",
    "Credencial pessoal e intransferível — nenhum profissional entra sem ela visível",
], Inches(0.75), Inches(3.25), Inches(12), size=13)

txb(s, "Escala por turno (preencher com nomes):", Inches(0.55), Inches(4.55), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["#", "Montagem 08/set", "Dia 1 — Manhã", "Dia 1 — Tarde", "Dia 2 — Manhã", "Dia 2 — Tarde"],
    [
        ("1–8",     "A definir", "A definir", "A definir", "A definir", "A definir"),
        ("9 (Fran)","—",         "Fran",      "—",         "Fran",      "—"),
    ],
    Inches(0.55), Inches(4.9),
    [Inches(0.85), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)],
    row_h=Inches(0.38),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 7 — ATIVAÇÃO: QUIZ BDC
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Ativação do Estande — Quiz BDC")

# Dois cards superiores
box(s, Inches(0.55), Inches(1.15), Inches(6.1), Inches(0.95), NAVY)
txb(s, "Totem touchscreen 43\" — Quiz interativo BDC",
    Inches(0.7), Inches(1.22), Inches(5.9), Inches(0.42),
    size=14, bold=True, color=AZUL_VIVO)
txb(s, "Visitante descobre qual solução SAP resolve sua dor em ~3 min",
    Inches(0.7), Inches(1.62), Inches(5.9), Inches(0.4),
    size=11, color=BRANCO)

box(s, Inches(7.0), Inches(1.15), Inches(6.0), Inches(0.95), NAVY)
txb(s, "Fornecedor: Entretec  ✓ aprovado 2025",
    Inches(7.15), Inches(1.22), Inches(5.8), Inches(0.42),
    size=14, bold=True, color=RGBColor(0x7F, 0xFF, 0x96))
txb(s, "Ref.: R$ 9.800 totem + R$ 500 frete + R$ 450/dia promotor (opcional)",
    Inches(7.15), Inches(1.62), Inches(5.8), Inches(0.4),
    size=11, color=BRANCO)

txb(s, "6 Pares Dor × Solução (foco BDC 2026):",
    Inches(0.55), Inches(2.3), Inches(12), Inches(0.38),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Dor (vermelho — visitante escolhe)", "Solução (verde — resultado)"],
    [
        ("Fechamento manual demorado",          "SAP BDC + Group Reporting"),
        ("Dados desconectados entre sistemas",  "SAP Datasphere"),
        ("Falta de visibilidade em tempo real", "SAP Analytics Cloud"),
        ("Planejamento orçamentário no Excel",  "SAP SAC Planning"),
        ("Integrações travando processos",       "SAP BTP"),
        ("Gestão fiscal engessada",              "SAP PaPM"),
    ],
    Inches(0.55), Inches(2.7),
    [Inches(6.5), Inches(6.15)],
    row_h=Inches(0.42),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — JORNADA DO VISITANTE
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Jornada do Visitante — Ativação Quiz")

steps = [
    ("1. Atração",        "Totem animado chama atenção. Equipe aborda:\n'Faz o quiz e descobre qual SAP resolve sua dor. Leva brinde!'"),
    ("2. Engajamento",    "Preenche nome, cargo, empresa, e-mail, WhatsApp.\nParticipa em menos de 3 minutos."),
    ("3. Qualificação",   "Quiz revela dor principal + solução SAP.\nLead categorizado: Quente / Morno / Frio."),
    ("4. Brinde + Conversa", "Equipe entrega brinde e inicia conversa consultiva:\n'Vi que você marcou [dor X]… posso mostrar como a Klabin resolveu isso?'"),
    ("5. CRM",            "Dados capturados vão para HubSpot.\nTag: SAP-NOW-2026 + classificação automática."),
    ("6. Follow-up D+3",  "SDR contata em até 3 dias:\n'Conforme seu quiz, acredito que podemos ajudar com [dor]. Agendamos?'"),
]

cols = 3
for i, (titulo, desc) in enumerate(steps):
    col = i % cols
    row = i // cols
    x = Inches(0.35 + col * 4.32)
    y = Inches(1.15 + row * 2.5)
    box(s, x, y, Inches(4.1), Inches(2.35), BG_LIGHT)
    # linha superior colorida
    box(s, x, y, Inches(4.1), Inches(0.04),
        AZUL_VIVO if i % 2 == 0 else RGBColor(0x7F, 0xFF, 0x96))
    box(s, x, y, Inches(4.1), Inches(0.44), NAVY)
    txb(s, titulo, x + Inches(0.1), y + Inches(0.06),
        Inches(3.9), Inches(0.38), size=13, bold=True, color=BRANCO)
    txb(s, desc, x + Inches(0.1), y + Inches(0.5),
        Inches(3.9), Inches(1.75), size=11, color=NAVY)

footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 9 — BRINDES
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Brindes")
txb(s, "Fornecedor referência: Unity Brindes (aprovado em 2025)",
    Inches(0.55), Inches(1.1), Inches(12), Inches(0.35),
    size=12, color=CINZA_TXT, italic=True)
table_block(s,
    ["Item", "Qtd", "Valor unit.", "Total"],
    [
        ("Bloco A5 capa dura (silk screen 1 cor)", "500", "~R$ 15,59", "~R$ 7.795"),
        ("Caneta esferográfica alumínio (laser)",  "500", "~R$ 5,37",  "~R$ 2.685"),
        ("**TOTAL",                               "",    "",           "**~R$ 10.480"),
    ],
    Inches(0.55), Inches(1.5),
    [Inches(7.5), Inches(1.5), Inches(2.1), Inches(2.1)],
    row_h=Inches(0.46),
)
txb(s, "Regras SAP para brindes:", Inches(0.55), Inches(3.35), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Devem ser sustentáveis — caneta + moleskine está aprovado pela SAP",
    "⚠️  Aprovação SAP: enviar junto com ativações até 31/07/2026",
    "Entrega prevista: até 14/agosto (fornecedor a confirmar)",
], Inches(0.75), Inches(3.75), Inches(12), size=14)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — ESTANDE
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Estande")

box(s, Inches(0.55), Inches(1.15), Inches(12.35), Inches(0.5), AMARELO)
txb(s, "⚠️  ATENÇÃO: Posição não privilegiada — toda estratégia deve ser de atração ativa.",
    Inches(0.7), Inches(1.2), Inches(12.0), Inches(0.42),
    size=13, bold=True, color=NAVY)

txb(s, "Layout do estande:", Inches(0.55), Inches(1.85), Inches(6.5), Inches(0.38),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Trainel (T1) — ~3,85m × 2,50m + espaço TV (1,45×0,83m) + balcão (0,50×1,00m) + Selo SAP Gold Partner",
    "Painel (P1) — Logo Solveplan vertical + BDC em destaque",
    "Balcão (B1) — Identidade visual + logo",
], Inches(0.75), Inches(2.25), Inches(6.3), size=13)

txb(s, "Conceito visual 2026:", Inches(7.1), Inches(1.85), Inches(5.8), Inches(0.38),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Foco em SAP Business Data Cloud",
    "Identidade Solveplan atualizada",
    "TV para vídeo looping",
    "Destaque: BDC | Datasphere | SAC | Group Reporting",
], Inches(7.3), Inches(2.25), Inches(5.6), size=13)

txb(s, "Cronograma de produção:", Inches(0.55), Inches(3.8), Inches(12), Inches(0.38),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Etapa", "Prazo"],
    [
        ("Orçamento com agência",        "Até 23/jun"),
        ("Arte final aprovada",           "Até 11/jul"),
        ("Pagamento 1ª parcela (50%)",    "10/jul"),
        ("Envio das artes exclusivas SAP","11/jul"),
        ("Pagamento 2ª parcela (50%)",    "Na entrega"),
    ],
    Inches(0.55), Inches(4.2),
    [Inches(9.0), Inches(3.6)],
    row_h=Inches(0.38),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 11 — VESTIMENTA, CAEX, EQUIPAMENTOS
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Vestimenta, Serviços CAEX e Equipamentos")

txb(s, "Vestimenta:", Inches(0.55), Inches(1.15), Inches(6.5), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Produto", "Qtd", "Valor unit.", "Total"],
    [("Tech T-Shirt personalizada", "25 un.", "~R$ 115", "~R$ 2.875")],
    Inches(0.55), Inches(1.5),
    [Inches(5.8), Inches(1.5), Inches(2.1), Inches(2.1)],
    row_h=Inches(0.42),
)

txb(s, "Serviços adicionais CAEX:", Inches(0.55), Inches(2.4), Inches(6.5), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Item", "Qtd", "Total"],
    [
        ("Coletor de dados (por dia)",             "2", "~R$ 1.100"),
        ("KVA energia (totem consume 2–3 KVAs)",   "2", "~R$ 780"),
        ("Internet (pacote a definir)",            "—", "A definir"),
    ],
    Inches(0.55), Inches(2.75),
    [Inches(7.5), Inches(1.5), Inches(2.55)],
    row_h=Inches(0.4),
)

txb(s, "Equipamentos:", Inches(7.3), Inches(1.15), Inches(5.8), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
equip = [
    "TV + suporte + cabos (verificar CAEX)",
    "1 Tablet para demos",
    "2 Pens drive vídeo looping (TV) + 1 backup",
    "2 Cabos HDMI + adaptador de rede",
    "Régua de energia / extensão",
    "Carregadores de celular + pano e álcool",
    "Agenda das sessões impressa",
    "2 displays QR Code (~R$ 22 un.) + 1 reserva",
]
bullets(s, equip, Inches(7.3), Inches(1.55), Inches(5.8), size=12, spacing=0.38)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 12 — CONTEÚDO DO EVENTO — KLABIN
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Conteúdo do Evento — Apresentação Principal")

box(s, Inches(0.55), Inches(1.15), Inches(12.35), Inches(0.52), AMARELO)
txb(s, "⚠️  Prazo de submissão para SAP: 26/06/2026 — NÃO PODE ATRASAR",
    Inches(0.7), Inches(1.2), Inches(12.0), Inches(0.42),
    size=13, bold=True, color=NAVY)

kv_block(s, [
    ("Tipo",          "Caso de Sucesso — Customer Story"),
    ("Formato",       "20 minutos | Auditório ~50 lugares | Sem Q&A"),
    ("Palestrante 1", "Executivo Klabin — CFO/CIO/Head de Dados (a confirmar até 10/06)"),
    ("Palestrante 2", "Alexandre Kuntgen (Solveplan)"),
], Inches(0.55), Inches(1.8), w1=Inches(2.8), w2=Inches(9.9))

txb(s, "Estrutura da sessão (20 min):", Inches(0.55), Inches(3.75), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Tempo", "Bloco", "Speaker"],
    [
        ("0–3 min",   "O desafio: como era antes, a dor real de dados na Klabin",  "Executivo Klabin"),
        ("3–8 min",   "A decisão: por que SAP BDC, por que Solveplan",            "Alexandre Kuntgen"),
        ("8–15 min",  "O que mudou: antes × depois com números reais",            "Executivo Klabin"),
        ("15–18 min", "O futuro: próximos passos com IA (alinha com tema do evento)", "Ambos"),
        ("18–20 min", "CTA: 'Venha ao nosso estande — estamos no slot [X]'",      "Alexandre Kuntgen"),
    ],
    Inches(0.55), Inches(4.12),
    [Inches(1.5), Inches(8.2), Inches(3.0)],
    row_h=Inches(0.42),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 13 — COMUNICAÇÃO: SOCIAL MEDIA & E-MAILS
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Comunicação — Social Media & E-mails")

txb(s, "Calendário de publicações (LinkedIn):",
    Inches(0.55), Inches(1.1), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Data", "Publicação"],
    [
        ("21/jul", "Primeiro post — 'Estaremos no SAP NOW 2026'"),
        ("30/jul", "Post de expectativa / sessão com Klabin anunciada"),
        ("11/ago", "Post de contagem regressiva"),
        ("01/set", "'Em 8 dias estaremos lá'"),
        ("08/set", "'É amanhã!' — com localização do estande"),
        ("09/set", "Cobertura ao vivo — Dia 1"),
        ("10/set", "Cobertura ao vivo — Dia 2"),
        ("11/set", "Post pós-evento — highlights + agradecimento"),
    ],
    Inches(0.55), Inches(1.48),
    [Inches(1.8), Inches(10.85)],
    row_h=Inches(0.38),
)

txb(s, "E-mails Marketing:", Inches(0.55), Inches(4.8), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Save the date → Base geral (Mai/Jun)",
    "Convite + link inscrição → Prospects e clientes (Jun/Jul)",
    "Lembrete sessão Klabin → Inscritos confirmados (Ago)  |  'É amanhã!' (08/set)  |  Agradecimento D+1",
], Inches(0.75), Inches(5.2), Inches(12.2), size=13)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 14 — ANÚNCIOS + UTMs
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Comunicação — Anúncios e UTMs")

# Card hotéis
box(s, Inches(0.55), Inches(1.15), Inches(5.85), Inches(2.6), NAVY)
txb(s, "1. Anúncios em Hotéis SP", Inches(0.7), Inches(1.22),
    Inches(5.6), Inches(0.45), size=15, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Período: 8, 9 e 10 de setembro",
    "~28 faces/telas em hotéis da região",
    "Valor ref.: R$ 2.805 (2025)",
    "Aprovação: antes de 31/ago",
], Inches(0.75), Inches(1.72), Inches(5.5), size=12,
    color=BRANCO, spacing=0.38)

# Card geolocalização
box(s, Inches(7.05), Inches(1.15), Inches(5.85), Inches(2.6), NAVY)
txb(s, "2. Geolocalização no Evento", Inches(7.2), Inches(1.22),
    Inches(5.6), Inches(0.45), size=15, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Google Ads + LinkedIn Ads",
    "Raio: 20km do Transamerica Expo Center",
    "Período: 8, 9 e 10 de setembro",
    "'Estamos no SAP NOW — venha ao estande Solveplan'",
    "Valor: a orçar",
], Inches(7.25), Inches(1.72), Inches(5.5), size=12,
    color=BRANCO, spacing=0.38)

txb(s, "UTMs Padrão:", Inches(0.55), Inches(3.95), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Canal", "UTM Source", "UTM Medium", "UTM Campaign"],
    [
        ("Organic Social",   "linkedin", "organic",    "sap-now-2026"),
        ("Paid Social",      "linkedin", "cpc",        "sap-now-2026"),
        ("E-mail Marketing", "email",    "newsletter", "sap-now-2026"),
        ("Form inscrição SAP","sap",     "event",      "sap-now-2026"),
    ],
    Inches(0.55), Inches(4.32),
    [Inches(3.8), Inches(2.5), Inches(2.5), Inches(3.85)],
    row_h=Inches(0.38),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 15 — GERAÇÃO DE DEMANDA SDR
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Geração de Demanda Pré-evento — SDR")

kv_block(s, [
    ("Objetivo",  "Preencher slots de diagnóstico BDC no estande antes do evento"),
    ("Meta",      "120 empresas abordadas → 12 agendas confirmadas"),
    ("Início",    "Junho 2026"),
    ("SDRs",      "2 SDRs (nomes a definir)"),
], Inches(0.55), Inches(1.15), w1=Inches(2.2), w2=Inches(10.5))

txb(s, "Processo:", Inches(0.55), Inches(3.2), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "1. SDR identifica decisor (CIO/CFO/Head de Dados) nas empresas-alvo",
    "2. Abordagem por e-mail + LinkedIn + telefone",
    "3. Convite para slot de diagnóstico BDC no estande (15–20 min)",
    "4. Confirmar agenda + registrar no CRM com tag SAP-NOW-2026",
], Inches(0.75), Inches(3.6), Inches(12.2), size=13)

box(s, Inches(0.55), Inches(5.35), Inches(12.35), Inches(0.82), NAVY)
txb(s, '📞  Script SDR:\n"[Nome], estaremos no SAP NOW AI Tour em setembro com demo exclusiva de SAP BDC. Tenho um slot de 15 min reservado para você no estande da Solveplan — posso garantir o seu?"',
    Inches(0.75), Inches(5.42), Inches(12.1), Inches(0.72),
    size=13, bold=True, color=RGBColor(0x7F, 0xFF, 0x96))
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 16 — CRONOGRAMA MACRO
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Cronograma Macro")
table_block(s,
    ["Data", "Entrega", "Responsável"],
    [
        ("18/05",    "Receber manual patrocinador CAEX",                    "Fran"),
        ("Mai",      "Lançar campanha 'Diagnóstico BDC Reservado'",         "Fran"),
        ("10/06",    "⚠️ Alinhar com Klabin — speaker e tema",              "Fran"),
        ("26/06",    "⚠️ Submeter sessão de conteúdo para SAP (CRÍTICO)",   "Fran"),
        ("02/jun",   "Liberação do site e abertura de registros",           "SAP"),
        ("Jun",      "Início abordagem SDR (2 SDRs)",                       "SDRs"),
        ("31/07",    "⚠️ Aprovação de ativações e brindes para SAP",        "Fran"),
        ("14/ago",   "Recebimento de brindes",                              "Fornecedor"),
        ("25/ago",   "App do evento lança / publicar slots diagnóstico",    "SAP / Fran"),
        ("Set/1",    "Treinamento e kickoff do time",                       "Fran"),
        ("08/set",   "Montagem do estande + entrega todos os itens",        "Time"),
        ("09–10/set","⭐ SAP NOW AI Tour Brazil 2026",                       "—"),
        ("10/set",   "Download mailing + início follow-up",                 "SDRs + Fran"),
        ("16/set",   "Relatório de resultados",                             "Fran"),
    ],
    Inches(0.4), Inches(1.15),
    [Inches(1.9), Inches(8.9), Inches(2.55)],
    row_h=Inches(0.37),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 17 — ORÇAMENTO
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Orçamento Previsto")
table_block(s,
    ["Item", "Valor previsto"],
    [
        ("Produção do estande (arte + layout)",            "~R$ 3.700"),
        ("Quiz/Totem interativo (Entretec)",               "~R$ 11.200"),
        ("Brindes (canetas + moleskine — 500 un.)",        "~R$ 10.480"),
        ("Camisetas equipe (25 un.)",                      "~R$ 2.875"),
        ("Cartões de visita",                              "~R$ 3.700"),
        ("Coletores de dados CAEX (2 un.)",                "~R$ 1.100"),
        ("KVA energia (2 un.)",                            "~R$ 780"),
        ("Anúncios hotéis SP",                             "~R$ 2.805"),
        ("Vídeo looping + depoimentos + teasers",          "~R$ 5.280"),
        ("Anúncios geolocalização (Google + LinkedIn)",    "A orçar"),
        ("Hotel equipe + Internet CAEX + Campanha LinkedIn","A definir"),
        ("Adicionais (displays, cabos, etc.)",             "~R$ 500"),
        ("**TOTAL DISPONÍVEL",                             "**R$ 276.510,21"),
    ],
    Inches(1.0), Inches(1.15),
    [Inches(9.8), Inches(2.8)],
    row_h=Inches(0.41),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 18 — CONVITE A CLIENTES
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Convite a Clientes C-Level")

box(s, Inches(0.55), Inches(1.15), Inches(12.35), Inches(0.52), NAVY)
txb(s, "Benefício Gold: convites ilimitados para clientes C-level do patrocinador",
    Inches(0.7), Inches(1.2), Inches(12.0), Inches(0.42),
    size=13, bold=True, color=RGBColor(0x7F, 0xFF, 0x96))

txb(s, "Critérios SAP:", Inches(0.55), Inches(1.9), Inches(6.5), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "C-level / Diretoria — preferencialmente em São Paulo",
    "Máximo 2 contatos por empresa/conta",
    "Todos passam por aprovação SAP antes da confirmação",
], Inches(0.75), Inches(2.3), Inches(6.0), size=13)

txb(s, "Processo:", Inches(7.1), Inches(1.9), Inches(5.8), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "1. Enviar convite personalizado",
    "2. Participante se inscreve pelo link",
    "3. SAP analisa e confirma",
    "4. Participante recebe e-mail de confirmação",
], Inches(7.3), Inches(2.3), Inches(5.6), size=13)

txb(s, "Lista estratégica — até 15 clientes (definir com sócios até 10/ago):",
    Inches(0.55), Inches(3.8), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Empresa", "Contato", "Cargo", "Status"],
    [("A definir", "—", "—", "Pendente")] * 5,
    Inches(0.55), Inches(4.18),
    [Inches(4.0), Inches(3.2), Inches(3.2), Inches(2.3)],
    row_h=Inches(0.4),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 19 — VÍDEO + FOLLOW-UP
# ════════════════════════════════════════════════════════════
s = new_slide()
title_block(s, "Vídeo Stand + Follow-up Pós-evento")

txb(s, "Vídeo looping (TV do estande):", Inches(0.55), Inches(1.15), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Cortes de cases (Klabin, clientes BDC) + comerciais SAP BDC + dashboards ao vivo",
    "Formato: MP4, 16:9, 1920×1080, looping contínuo  |  2 pens drive (principal + backup)",
    "Produção: a contratar  |  Prazo de entrega: até 25/ago",
], Inches(0.75), Inches(1.55), Inches(12.2), size=13)

txb(s, "Vídeos durante o evento:", Inches(0.55), Inches(2.75), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
bullets(s, [
    "Shorts 9:16 para redes sociais  |  Depoimentos curtos com clientes  |  Gravação de cases",
    "Pacote depoimentos 4 vídeos ~R$ 3.680  |  Teasers (4 un.) ~R$ 1.600  |  Edição adicional ~R$ 890/un.",
], Inches(0.75), Inches(3.15), Inches(12.2), size=13)

txb(s, "Follow-up pós-evento (estruturado):", Inches(0.55), Inches(4.05), Inches(12), Inches(0.35),
    size=13, bold=True, color=AZUL_VIVO)
table_block(s,
    ["Dia", "Ação"],
    [
        ("D+1", "Consolidar leads + limpar dados + subir no HubSpot com tag SAP-NOW-2026"),
        ("D+1", "Contato direto (e-mail/WhatsApp) — leads quentes têm prioridade"),
        ("D+3", "E-mail com material de valor (case BDC relevante para a dor do lead)"),
        ("D+5", "Ligação — confirmar interesse e propor agendamento"),
        ("D+7", "Último contato da sequência — agendar ou encaminhar para nutrição"),
    ],
    Inches(0.55), Inches(4.43),
    [Inches(1.3), Inches(11.35)],
    row_h=Inches(0.38),
)
footer(s)


# ════════════════════════════════════════════════════════════
# SLIDE 20 — PRÓXIMOS PASSOS (encerramento escuro)
# ════════════════════════════════════════════════════════════
s = new_slide(dark=True)
title_block(s, "Próximos Passos Críticos", dark=True)

table_block(s,
    ["Prazo", "Ação"],
    [
        ("Hoje 18/05", "Ler manual CAEX — confirmar posição e regras do estande 2026"),
        ("Mai",        "Lançar campanha 'Diagnóstico BDC Reservado' (LinkedIn Ads)"),
        ("10/06",      "⚠️ Alinhar com Klabin — speaker, cargo, tema da sessão"),
        ("26/06",      "⚠️ Submeter sessão de conteúdo para SAP  ← CRÍTICO"),
        ("31/07",      "⚠️ Aprovação de ativações + brindes para SAP  ← CRÍTICO"),
        ("08/set",     "Montagem do estande + entrega de todos os itens"),
        ("09–10/set",  "⭐ SAP NOW AI Tour Brazil 2026"),
    ],
    Inches(0.55), Inches(1.35),
    [Inches(2.3), Inches(10.4)],
    row_h=Inches(0.5),
    hdr_color=AZUL_VIVO,
)

txb(s, "SAP NOW AI Tour Brazil 2026  |  9–10 Set  |  Transamérica Expo Center  |  Meta: 130 leads · 28 reuniões · R$ 3,2M pipeline",
    Inches(0.55), Inches(6.25), Inches(12.35), Inches(0.42),
    size=11, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

footer(s, dark=True, page_label="")


# ── Salvar ────────────────────────────────────────────────────
output = (r"c:\Users\franc\solveplan.com\Roberto Molina - Marketing"
          r"\1. MKT Estrategy\3. Agentes de IA\ccos-ratos\eventos"
          r"\sap-now-2026\SAP-NOW-2026-Planejamento-Solveplan.pptx")
prs.save(output)
print(f"Salvo: {output}")
print(f"Total slides: {len(prs.slides)}")
