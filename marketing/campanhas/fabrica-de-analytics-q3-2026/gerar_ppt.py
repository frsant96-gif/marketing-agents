from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Cores Solveplan
AZUL        = RGBColor(0x00, 0x2F, 0x6C)
AMARELO     = RGBColor(0xF5, 0xA8, 0x00)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF4, 0xF4, 0xF4)
CINZA_TEXTO = RGBColor(0x44, 0x44, 0x44)
AZUL_CLARO  = RGBColor(0xE8, 0xF0, 0xFE)
VERMELHO    = RGBColor(0xC0, 0x39, 0x2B)
VERDE       = RGBColor(0x1A, 0x7A, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank_slide():
    return prs.slide_layouts[6]

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=12, bold=False, color=RGBColor(0,0,0),
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def divider(slide, l, t, w, color=AMARELO):
    rect(slide, l, t, w, 0.05, color)

def header(slide, title):
    rect(slide, 0, 0, 13.33, 0.12, AMARELO)
    rect(slide, 0, 6.8, 13.33, 0.12, AZUL)
    txt(slide, title, 0.5, 0.25, 12, 0.65, 22, bold=True, color=AZUL)
    divider(slide, 0.5, 0.95, 12.3)

# ─────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
rect(slide, 0, 0, 13.33, 7.5, AZUL)
rect(slide, 0, 6.8, 13.33, 0.7, AMARELO)
rect(slide, 0, 0, 13.33, 0.12, AMARELO)

txt(slide, "Campanha de Marketing", 0.8, 1.4, 10, 0.8, 16, color=RGBColor(0xCC,0xCC,0xCC))
txt(slide, "Fábrica de Analytics", 0.8, 2.1, 11, 1.1, 38, bold=True, color=BRANCO)
txt(slide, "Geração de Leads — Q3 2026", 0.8, 3.2, 11, 0.7, 22, color=AMARELO)
txt(slide, '"A implantação do SAP foi só o começo. Quem cuida do que vem depois?"', 0.8, 3.95, 11, 0.6, 13, italic=True, color=RGBColor(0xCC,0xCC,0xCC))
txt(slide, "Julho · Agosto · Setembro 2026  |  Suporte e melhorias em SAC / Datasphere por demanda", 0.8, 4.6, 11, 0.5, 11, color=RGBColor(0x88,0x88,0x88))
txt(slide, "solveplan", 0.8, 6.88, 4, 0.45, 14, bold=True, color=AZUL)

# ─────────────────────────────────────────
# SLIDE 2 — DORES REAIS (VOC)
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
header(slide, "As Dores Reais dos Clientes — VoC de 18 Reuniões")

txt(slide, "Baseado em falas explícitas dos clientes. Nenhuma inferência.", 0.5, 1.0, 12, 0.35, 10, italic=True, color=CINZA_TEXTO)

dores = [
    ("Sem parceiro\npara suportar SAC",
     '"a gente queria uma empresa para dar suporte\nnessa passagem de conhecimento"',
     "DOR CENTRAL", VERDE),
    ("SAC contratado,\nsem suporte p/ ativar",
     '"tentamos usar, não foi muito bem aceito"\n"temos licenciamento... não chegamos a evoluir"',
     "4 reuniões", AMARELO),
    ("RFC em risco\nde bloqueio SAP",
     '"nós estamos assombrados com isso"\n"sem ele basicamente a fábrica pararia"',
     "3 reuniões", VERMELHO),
    ("S4 sobrecarregado\nprecisa ser otimizado",
     '"eu preciso desafogar o nosso S4"\n"vai vir uma fatura no final"',
     "3 reuniões", AZUL),
]

x = 0.4
for titulo, fala, freq, cor in dores:
    rect(slide, x, 1.45, 3.0, 4.8, AZUL_CLARO)
    rect(slide, x, 1.45, 3.0, 0.08, cor)
    txt(slide, titulo, x+0.12, 1.55, 2.76, 0.75, 12, bold=True, color=AZUL)
    divider(slide, x+0.12, 2.32, 2.76, cor)
    txt(slide, fala, x+0.12, 2.42, 2.76, 2.5, 10, italic=True, color=CINZA_TEXTO)
    rect(slide, x+0.12, 5.4, 2.76, 0.5, cor)
    txt(slide, freq, x+0.12, 5.4, 2.76, 0.5, 13, bold=True, color=BRANCO if cor != AMARELO else AZUL, align=PP_ALIGN.CENTER)
    x += 3.22

txt(slide, "⚠  O que NÃO usar: IA/preditivo, BDC como posicionamento amplo, Joule/Copilot — zero menções espontâneas nas 18 reuniões.",
    0.4, 6.3, 12.5, 0.4, 10, italic=True, color=CINZA_TEXTO)

# ─────────────────────────────────────────
# SLIDE 3 — DOR RAIZ UNIFICADA
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
header(slide, "O Padrão que Atravessa Todas as Dores")

rect(slide, 0.4, 1.1, 12.5, 1.1, AZUL)
txt(slide, '"A implantação do SAP foi só o começo. Quem cuida do que vem depois?"',
    0.7, 1.2, 12.0, 0.8, 18, bold=True, italic=True, color=AMARELO, align=PP_ALIGN.CENTER)

txt(slide, "5 dores diferentes nas 18 reuniões — mesma raiz:",
    0.5, 2.4, 12.3, 0.4, 12, bold=True, color=AZUL)

raizes = [
    ("SAC parado",          "Ferramenta implantada, sem suporte para ativar e manter",           "4 reuniões"),
    ("Sem parceiro técnico","Time de TI sem capacidade para evoluir o ambiente",                  "3+ reuniões"),
    ("RFC em risco",        "Arquitetura crítica sem suporte para migrar antes do bloqueio SAP",  "3 reuniões"),
    ("S4 sobrecarregado",   "Relatórios no transacional sem suporte para mover",                  "3 reuniões"),
    ("SAC só como BI",      "Ambiente funciona, mas sem apoio para evoluir para planning",        "3 reuniões"),
]

y = 2.9
for dor, raiz, freq in raizes:
    rect(slide, 0.4, y, 12.5, 0.6, AZUL_CLARO if raizes.index((dor,raiz,freq)) % 2 == 0 else BRANCO)
    rect(slide, 0.4, y, 0.08, 0.6, AMARELO)
    txt(slide, dor,  0.6, y+0.1, 2.5, 0.4, 11, bold=True, color=AZUL)
    txt(slide, raiz, 3.2, y+0.1, 7.5, 0.4, 11, color=CINZA_TEXTO)
    txt(slide, freq, 10.8, y+0.1, 2.0, 0.4, 10, italic=True, color=CINZA_TEXTO, align=PP_ALIGN.RIGHT)
    y += 0.62

rect(slide, 0.4, y+0.05, 12.5, 0.55, AZUL)
txt(slide, "Dor comum: têm a ferramenta SAP implantada, mas não têm capacidade técnica interna para ativar, sustentar, melhorar e evoluir.",
    0.6, y+0.1, 12.1, 0.42, 11, bold=True, color=BRANCO)

# ─────────────────────────────────────────
# SLIDE 4 — PÚBLICO-ALVO E ICP
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
header(slide, "Público-Alvo")

# Card esquerda — perfil
rect(slide, 0.4, 1.1, 5.9, 5.3, AZUL_CLARO)
txt(slide, "Perfil do ICP", 0.55, 1.2, 5.6, 0.45, 14, bold=True, color=AZUL)
divider(slide, 0.55, 1.68, 5.6)

campos = [
    ("Cargos",      "Gerente/Coord. de TI, Head de Dados/BI\nGestor de FP&A, Controller, CFO"),
    ("Empresa",     "Médio e grande porte\nS4HANA ou RISE com SAC contratado"),
    ("Segmentos",   "Indústria · Varejo · Agronegócio · Financeiro"),
    ("Tamanho",     "Faturamento > R$ 500M"),
]
y = 1.78
for label, valor in campos:
    txt(slide, label, 0.55, y, 2.0, 0.55, 11, bold=True, color=AZUL)
    txt(slide, valor, 2.5, y, 3.7, 0.55, 11, color=CINZA_TEXTO)
    y += 0.65

# Card direita — momento de compra
rect(slide, 6.8, 1.1, 6.1, 5.3, AZUL_CLARO)
txt(slide, "Momento de Compra", 6.95, 1.2, 5.9, 0.45, 14, bold=True, color=AZUL)
divider(slide, 6.95, 1.68, 5.9)

momentos = [
    ("SAC parado",      "Licença ativa, ferramenta sem uso"),
    ("RFC em risco",    "Arquitetura depende de RFC — nota SAP vigente"),
    ("S4 pesado",       "Relatórios rodando no transacional"),
    ("Sem time",        "TI sem capacidade técnica para evoluir SAC/Datasphere"),
    ("Pós-RISE",        "Go Live recente, jornada analítica não iniciada"),
]
y = 1.78
for gatilho, desc in momentos:
    rect(slide, 6.95, y, 5.9, 0.62, AZUL)
    txt(slide, gatilho, 7.1, y+0.08, 1.8, 0.44, 11, bold=True, color=AMARELO)
    txt(slide, desc, 8.95, y+0.08, 3.8, 0.44, 11, color=BRANCO)
    y += 0.72

# ─────────────────────────────────────────
# SLIDE 5 — 5 ÂNGULOS DE MENSAGEM
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
header(slide, "5 Ângulos de Mensagem — Mesma Raiz, Dores Diferentes")

txt(slide, "Mensagem guarda-chuva: \"A implantação do SAP foi só o começo. Quem cuida do que vem depois?\"",
    0.5, 1.0, 12.3, 0.38, 11, bold=True, italic=True, color=AZUL)

angulos = [
    ("01", "Sem parceiro\ntécnico",
     "Quem cuida do seu ambiente SAC depois da implantação? Banco de horas com especialistas — você chama quando precisa.",
     VERDE),
    ("02", "SAC parado,\nsem suporte",
     "Não é falta de ferramenta — é falta de suporte para ativar e manter. Coloca no ar sem projeto de meses.",
     AMARELO),
    ("03", "RFC\nurgente",
     "A SAP vai bloquear extrações via RFC. Com suporte especializado, existe rota de transição sem reescrever tudo.",
     VERMELHO),
    ("04", "S4\nsobrecarregado",
     "Seu S4 não é para relatórios. Com suporte técnico certo, você move para a camada analítica sem projeto interno.",
     AZUL),
    ("05", "SAC só como\nBI — melhorias",
     "Seu SAC faz muito mais que dashboards. Planning, forecast, simulações — a Solveplan evolui por demanda.",
     RGBColor(0x6A, 0x0D, 0x83)),
]

# layout: 3 em cima + 2 em baixo centralizados
positions_5 = [
    (0.25, 1.5), (4.55, 1.5), (8.85, 1.5),
    (2.4,  4.2), (6.7,  4.2),
]
for i, (num, titulo, msg, cor) in enumerate(angulos):
    x, y = positions_5[i]
    rect(slide, x, y, 4.05, 2.55, AZUL_CLARO)
    rect(slide, x, y, 0.5, 2.55, cor)
    txt(slide, num, x, y+1.0, 0.5, 0.5, 13, bold=True, color=BRANCO if cor != AMARELO else AZUL, align=PP_ALIGN.CENTER)
    txt(slide, titulo, x+0.6, y+0.08, 3.35, 0.55, 11, bold=True, color=AZUL)
    divider(slide, x+0.6, y+0.65, 3.35, cor)
    txt(slide, f'"{msg}"', x+0.6, y+0.75, 3.35, 1.6, 9, italic=True, color=CINZA_TEXTO)

# ─────────────────────────────────────────
# SLIDE 5 — METAS Q3 2026
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
header(slide, "Metas de Sucesso — Q3 2026")

txt(slide, "Conservadoras e defensáveis para um nicho B2B SAP. 1 deal fechado paga o investimento inteiro da campanha.",
    0.5, 1.0, 12.3, 0.38, 10, italic=True, color=CINZA_TEXTO)

metas = [
    ("30–50 mil",    "Alcance\n(impressões)"),
    ("8–12",         "Leads\nGerados"),
    ("4–7",          "MQLs\n(qualificados)"),
    ("2–4",          "Reuniões\nAgendadas"),
    ("1–2",          "Oportunidades\nAbertas"),
    ("R$ 678k\n–1,3M", "Pipeline\nGerado"),
]

positions_m = [
    (0.4,  1.5), (2.65, 1.5), (4.9,  1.5),
    (7.15, 1.5), (9.4,  1.5), (11.1, 1.5),
]
widths_m = [2.0, 2.0, 2.0, 2.0, 2.0, 2.0]
for i, (num, label) in enumerate(metas):
    x, y = positions_m[i]
    rect(slide, x, y, 1.9, 3.8, AZUL)
    txt(slide, num, x, y+0.5, 1.9, 1.5, 20, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y+2.3, 1.9, 1.0, 10, color=BRANCO, align=PP_ALIGN.CENTER)

txt(slide, "Ticket médio Solveplan: R$ 678.406  |  Meta de CPL: < R$ 2.500  |  MQL/Lead: > 50%",
    0.5, 5.6, 12.3, 0.4, 11, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

rect(slide, 0.4, 6.1, 12.5, 0.55, AZUL_CLARO)
txt(slide, "Racional: com 8–12 leads qualificados e ticket médio de R$ 678k, 1 deal fecha R$ 678k de pipeline — ROI positivo em qualquer cenário de budget razoável.",
    0.55, 6.13, 12.2, 0.45, 10, italic=True, color=AZUL)

# ─────────────────────────────────────────
# SLIDE 6 — CANAIS
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
header(slide, "Canais e Formatos")

canais = [
    ("Email\nMarketing", "Reativar 95 warm leads\ndos eventos H1 2026", "Alta", "Custo zero — base já qualificada", AMARELO),
    ("LinkedIn\nOrgânico", "4 posts (um por ângulo\nde dor)", "Alta", "ICP ativo no LinkedIn", AZUL),
    ("LinkedIn\nAds", "Lead gen form +\nsponsored content", "Alta", "Segmentação precisa por cargo + empresa SAP", AZUL),
    ("Google\nAds", "Keywords RFC,\nSAC Planning, Datasphere", "Alta", "CPC R$ 2,48 | CTR 4,97% — já performa", VERDE),
    ("Blog /\nArtigo", "1 artigo por ângulo\n(SEO)", "Média", "Médio prazo — autoridade e tráfego orgânico", CINZA_TEXTO),
    ("WhatsApp\nDireto", "Follow-up dos\n95 warm leads", "URGENTE", "Primeira ação — antes de qualquer mídia paga", VERMELHO),
]

x = 0.35
for canal, formato, prioridade, nota, cor in canais:
    rect(slide, x, 1.1, 2.05, 5.3, AZUL_CLARO)
    rect(slide, x, 1.1, 2.05, 0.08, cor)
    txt(slide, canal, x+0.1, 1.2, 1.85, 0.65, 12, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    divider(slide, x+0.1, 1.87, 1.85, cor)
    txt(slide, formato, x+0.1, 1.97, 1.85, 1.3, 10, color=CINZA_TEXTO)
    cor_pri = VERMELHO if prioridade == "URGENTE" else (AMARELO if prioridade == "Alta" else CINZA_TEXTO)
    rect(slide, x+0.1, 3.35, 1.85, 0.42, cor_pri)
    txt(slide, prioridade, x+0.1, 3.35, 1.85, 0.42, 11, bold=True,
        color=AZUL if prioridade == "Alta" or prioridade == "URGENTE" else BRANCO,
        align=PP_ALIGN.CENTER)
    txt(slide, nota, x+0.1, 3.88, 1.85, 1.3, 9, italic=True, color=CINZA_TEXTO)
    x += 2.18

# ─────────────────────────────────────────
# SLIDE 7 — CRONOGRAMA
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
header(slide, "Cronograma — Q3 2026")

fases = [
    ("PRÉ-CAMPANHA\nAté 04/07", [
        "Segmentar 95 warm leads",
        "Criar sequência de emails",
        "Configurar LinkedIn Ads + Google Ads",
        "Criar landing page",
        "Briefar time comercial",
    ], CINZA_TEXTO),
    ("JULHO\n07–31/07", [
        "Semana 1: Reativar warm leads — RFC",
        "Semana 2: LinkedIn posts + Ads ativos",
        "Semana 3: Email RFC + Google Ads",
        "Semana 4: Artigo RFC + posts ângulo 2",
    ], AZUL),
    ("AGOSTO\n04–29/08", [
        "Semana 5: Email S4 sobrecarregado",
        "Semana 6: Posts ângulo capacidade técnica",
        "Semana 7: Email SAC como BI + artigo",
        "Semana 8: Revisão mid-campaign",
    ], AZUL),
    ("SETEMBRO\n01–30/09", [
        "Acelerar canal líder (definir em ago)",
        "Reativar leads não convertidos",
        "Follow-up e qualificação comercial",
        "Relatório Q3 + aprendizados para Q4",
    ], VERDE),
]

x = 0.4
for mes, atividades, cor in fases:
    rect(slide, x, 1.1, 3.0, 0.65, cor if cor != CINZA_TEXTO else RGBColor(0x88,0x88,0x88))
    txt(slide, mes, x, 1.1, 3.0, 0.65, 11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    rect(slide, x, 1.75, 3.0, 4.65, AZUL_CLARO)
    y = 1.85
    for at in atividades:
        txt(slide, f"• {at}", x+0.12, y, 2.78, 0.55, 10, color=CINZA_TEXTO)
        y += 0.58
    x += 3.22

# ─────────────────────────────────────────
# SLIDE 8 — BUDGET E PEÇAS
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
header(slide, "Budget Estimado e Peças Necessárias")

# Budget
txt(slide, "Budget estimado", 0.5, 1.1, 6.0, 0.45, 14, bold=True, color=AZUL)
divider(slide, 0.5, 1.57, 6.0)

itens_budget = [
    ("LinkedIn Ads",          "R$ 3.000–6.000/mês"),
    ("Google Ads",            "R$ 1.500–3.000/mês"),
    ("Produção de peças",     "R$ 1.500–2.500 (único)"),
    ("Total Q3 estimado",     "R$ 16.000–29.000"),
]
y = 1.65
for i, (item, valor) in enumerate(itens_budget):
    is_total = "Total" in item
    cor_bg = AZUL if is_total else (CINZA_CLARO if i % 2 == 0 else BRANCO)
    rect(slide, 0.5, y, 6.0, 0.5, cor_bg)
    cor_txt = BRANCO if is_total else CINZA_TEXTO
    txt(slide, item,  0.6, y+0.07, 3.5, 0.36, 11, bold=is_total, color=cor_txt)
    txt(slide, valor, 4.2, y+0.07, 2.2, 0.36, 11, bold=is_total, color=AMARELO if is_total else CINZA_TEXTO)
    y += 0.52

rect(slide, 0.5, y+0.1, 6.0, 0.55, AZUL_CLARO)
txt(slide, "Google Ads atual: CPC R$ 2,48 · CTR 4,97% — canal de menor risco para escalar.",
    0.6, y+0.15, 5.8, 0.42, 10, italic=True, color=AZUL)

# Peças
txt(slide, "Peças necessárias", 7.2, 1.1, 5.7, 0.45, 14, bold=True, color=AZUL)
divider(slide, 7.2, 1.57, 5.7)

pecas = [
    ("Email sem parceiro (3 emails)", "Email marketing",   "/material-campanha"),
    ("Email RFC urgência (3 emails)", "Email marketing",   "/material-campanha"),
    ("Email SAC parado (3 emails)",   "Email marketing",   "/material-campanha"),
    ("4 posts LinkedIn",              "LinkedIn orgânico", "/post-social"),
    ("Banners LinkedIn Ads (3 var)",  "LinkedIn Ads",      "/anuncio"),
    ("Landing page suporte SAC",      "Paid + orgânico",   "/material-campanha"),
    ("Artigo RFC bloqueado",          "Blog / SEO",        "/artigo-blog"),
    ("Artigo quem sustenta seu SAC",  "Blog / SEO",        "/artigo-blog"),
    ("Copy Google Ads",               "Google Ads",        "/anuncio"),
]
y = 1.65
headers_p = ["Peça", "Canal", "Skill"]
col_x = [7.2, 9.8, 11.6]
col_w = [2.5, 1.7, 1.5]
rect(slide, 7.2, y, 6.0, 0.42, AMARELO)
for j, h in enumerate(headers_p):
    txt(slide, h, col_x[j]+0.1, y+0.05, col_w[j], 0.32, 10, bold=True, color=AZUL)
y += 0.44
for i, (peca, canal, skill) in enumerate(pecas):
    cor_bg = CINZA_CLARO if i % 2 == 0 else BRANCO
    rect(slide, 7.2, y, 6.0, 0.42, cor_bg)
    txt(slide, peca,  col_x[0]+0.1, y+0.05, col_w[0], 0.32, 9, color=CINZA_TEXTO)
    txt(slide, canal, col_x[1]+0.1, y+0.05, col_w[1], 0.32, 9, color=CINZA_TEXTO)
    txt(slide, skill, col_x[2]+0.1, y+0.05, col_w[2], 0.32, 9, bold=True, color=AZUL)
    y += 0.44

# ─────────────────────────────────────────
# SLIDE 9 — PRÓXIMOS PASSOS
# ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_slide())
rect(slide, 0, 0, 13.33, 7.5, AZUL)
rect(slide, 0, 0, 13.33, 0.12, AMARELO)
rect(slide, 0, 6.8, 13.33, 0.7, AMARELO)

txt(slide, "Próximos Passos", 0.8, 0.35, 11, 0.7, 28, bold=True, color=BRANCO)
divider(slide, 0.8, 1.1, 11.7, AMARELO)

passos = [
    ("1", "Reativar os 95 warm leads dos eventos H1 via WhatsApp/email — ação imediata, custo zero", "Esta semana"),
    ("2", "Segmentar base por dor (RFC / SAC parado / S4 heavy) no HubSpot",                         "Até 27/06"),
    ("3", "Criar sequência de emails por ângulo de dor — /material-campanha",                         "Até 04/07"),
    ("4", "Configurar LinkedIn Ads + Google Ads com ângulos RFC e SAC parado",                        "Até 04/07"),
    ("5", "Criar landing page com mensagem guarda-chuva: 'a implantação foi só o começo'",            "Até 04/07"),
    ("6", "Publicar artigo: 'RFC bloqueado — o que muda no seu ambiente SAP'",                        "Semana 4/jul"),
    ("7", "Revisar performance mid-campaign e otimizar canais",                                       "25/08"),
    ("8", "Relatório Q3 + planejamento Q4 com os aprendizados",                                       "07/10"),
]

y = 1.25
for num, acao, prazo in passos:
    rect(slide, 0.8, y, 0.6, 0.52, AMARELO)
    txt(slide, num, 0.8, y, 0.6, 0.52, 14, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    txt(slide, acao,  1.55, y+0.07, 9.2, 0.38, 11, color=BRANCO)
    txt(slide, prazo, 10.85, y+0.07, 2.1, 0.38, 11, color=AMARELO, align=PP_ALIGN.RIGHT)
    y += 0.62

txt(slide, "solveplan", 0.8, 6.88, 4, 0.45, 14, bold=True, color=AZUL)

# ─────────────────────────────────────────
# SALVAR
# ─────────────────────────────────────────
output = (
    r"c:\Users\franc\solveplan.com\Roberto Molina - Marketing"
    r"\1. MKT Estrategy\3. Agentes de IA\ccos-ratos"
    r"\marketing\campanhas\fabrica-de-analytics-q3-2026"
    r"\Fabrica_Analytics_Q3_2026_Campanha.pptx"
)
prs.save(output)
print(f"PPT salvo: {output}")
