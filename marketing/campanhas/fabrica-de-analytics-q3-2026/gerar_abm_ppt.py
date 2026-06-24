from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL        = RGBColor(0x00, 0x2F, 0x6C)
AMARELO     = RGBColor(0xF5, 0xA8, 0x00)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF4, 0xF4, 0xF4)
CINZA_TEXTO = RGBColor(0x44, 0x44, 0x44)
AZUL_CLARO  = RGBColor(0xE8, 0xF0, 0xFE)
VERMELHO    = RGBColor(0xC0, 0x39, 0x2B)
VERDE       = RGBColor(0x1A, 0x7A, 0x4A)
ROXO        = RGBColor(0x6A, 0x0D, 0x83)
LARANJA     = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank():
    return prs.slide_layouts[6]

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=12, bold=False, color=RGBColor(0,0,0),
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def div(slide, l, t, w, color=AMARELO):
    rect(slide, l, t, w, 0.05, color)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.12, AMARELO)
    rect(slide, 0, 6.88, 13.33, 0.12, AZUL)
    txt(slide, title, 0.5, 0.2, 12, 0.6, 20, bold=True, color=AZUL)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.75, 12, 0.3, 10, italic=True, color=CINZA_TEXTO)
    div(slide, 0.5, 0.95 if subtitle else 0.85, 12.3)

# ═══════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
rect(slide, 0, 0, 13.33, 7.5, AZUL)
rect(slide, 0, 0, 13.33, 0.12, AMARELO)
rect(slide, 0, 6.88, 13.33, 0.62, AMARELO)
rect(slide, 6.5, 0.12, 0.05, 6.76, AMARELO)

txt(slide, "Campanha Integrada — Q3 2026", 0.8, 1.2, 5.5, 0.55, 14, color=RGBColor(0xCC,0xCC,0xCC))
txt(slide, "Fábrica de\nAnalytics", 0.8, 1.8, 5.5, 1.5, 36, bold=True, color=BRANCO)
txt(slide, "Geração de Demanda + ABM", 0.8, 3.35, 5.5, 0.6, 18, color=AMARELO, bold=True)
txt(slide, '"A implantação do SAP foi só o começo.\nQuem cuida do que vem depois?"',
    0.8, 4.05, 5.5, 0.9, 12, italic=True, color=RGBColor(0xBB,0xBB,0xBB))
txt(slide, "Julho · Agosto · Setembro 2026", 0.8, 5.1, 5.5, 0.4, 11, color=RGBColor(0x88,0x88,0x88))

# Lado direito — números
stats = [
    ("156", "empresas mapeadas"),
    ("7",   "contas Tier 1 (VoC)"),
    ("95",  "warm leads para reativar"),
    ("2+",  "oportunidades esperadas"),
]
y = 1.4
for num, label in stats:
    rect(slide, 7.2, y, 5.6, 0.95, RGBColor(0x00, 0x22, 0x55))
    txt(slide, num,   7.3, y+0.05, 1.5, 0.75, 28, bold=True, color=AMARELO)
    txt(slide, label, 8.9, y+0.22, 3.7, 0.5,  12, color=BRANCO)
    y += 1.1

txt(slide, "solveplan", 0.8, 6.96, 4, 0.38, 13, bold=True, color=AZUL)

# ═══════════════════════════════════════
# SLIDE 2 — ESTRATÉGIA: 2 MOVIMENTOS
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
header(slide, "Dois Movimentos Simultâneos — Mesma Mensagem")

txt(slide, '"A implantação do SAP foi só o começo. Quem cuida do que vem depois?"',
    0.5, 1.0, 12.3, 0.38, 12, bold=True, italic=True, color=AZUL, align=PP_ALIGN.CENTER)

# Card esquerdo — Demand Gen
rect(slide, 0.4, 1.5, 5.9, 5.1, AZUL_CLARO)
rect(slide, 0.4, 1.5, 5.9, 0.08, AMARELO)
txt(slide, "Geração de Demanda", 0.55, 1.58, 5.6, 0.5, 14, bold=True, color=AZUL)
txt(slide, "Para quem ainda não nos conhece", 0.55, 2.08, 5.6, 0.3, 10, italic=True, color=CINZA_TEXTO)
div(slide, 0.55, 2.42, 5.6, AMARELO)
itens_dg = [
    "Email marketing — 95 warm leads (reativação imediata)",
    "LinkedIn Ads — ICP geral (fora da lista ABM)",
    "Google Ads — keywords RFC + SAC + Datasphere",
    "LinkedIn orgânico — 4-5 posts/mês",
    "Blog — 2 artigos SEO (RFC + quem sustenta SAC)",
]
y = 2.55
for item in itens_dg:
    txt(slide, f"• {item}", 0.55, y, 5.6, 0.42, 10, color=CINZA_TEXTO)
    y += 0.45
rect(slide, 0.55, 5.55, 5.6, 0.6, AZUL)
txt(slide, "Meta: 8–12 leads | 2–4 reuniões | R$ 678k–1,3M pipeline",
    0.7, 5.6, 5.3, 0.48, 10, bold=True, color=AMARELO)

# Card direito — ABM
rect(slide, 7.0, 1.5, 5.9, 5.1, AZUL_CLARO)
rect(slide, 7.0, 1.5, 5.9, 0.08, VERDE)
txt(slide, "Account-Based Marketing", 7.15, 1.58, 5.6, 0.5, 14, bold=True, color=AZUL)
txt(slide, "Para quem já foi impactado nos eventos H1 2026", 7.15, 2.08, 5.6, 0.3, 10, italic=True, color=CINZA_TEXTO)
div(slide, 7.15, 2.42, 5.6, VERDE)
itens_abm = [
    "Tier 1: 7 contas VoC — InMail personalizado + SDR",
    "Tier 2: ~50 contas por vertical (Agro, Energia, Indústria)",
    "Tier 3: ~99 contas — matched audience LinkedIn",
    "Orquestração com vendas — alerta de engajamento",
    "Mensuração por conta — cobertura e pipeline ABM",
]
y = 2.55
for item in itens_abm:
    txt(slide, f"• {item}", 7.15, y, 5.6, 0.42, 10, color=CINZA_TEXTO)
    y += 0.45
rect(slide, 7.15, 5.55, 5.6, 0.6, VERDE)
txt(slide, "Meta: 2–3 reuniões Tier 1 | 1–2 OPPs ABM | R$ 678k–2M pipeline",
    7.3, 5.6, 5.3, 0.48, 10, bold=True, color=BRANCO)

# ═══════════════════════════════════════
# SLIDE 3 — TIERS ABM
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
header(slide, "Segmentação ABM — 156 Empresas em 3 Tiers", "Base: eventos SAP BTP Experience + IT Summit Agro 2026")

tiers = [
    ("TIER 1", "1:1 Personalizado", "7 contas",
     "Contas do VoC com dor\nexplícita identificada",
     "InMail personalizado\n+ SDR outreach\n+ Sponsored Content",
     "R$ 250/conta/mês",
     VERMELHO,
     "COMIGO · Geoavia · Tuberfil\nGrupo Soufer · Docol\nAlgar · BTG Pactual"),
    ("TIER 2", "1:Few por Vertical", "~50 contas",
     "Líderes de TI/Dados dos\neventos — agrupados por setor",
     "Sponsored Content A/B\npor vertical\n+ SDR reforçado",
     "R$ 60/conta/mês",
     AMARELO,
     "Agro/Bioenergia · Energia/Utilities\nIndústria/Manufatura"),
    ("TIER 3", "1:Many ICP", "~99 contas",
     "Demais empresas dos eventos\ncom perfil SAP",
     "Matched Audience LinkedIn\nMensagem guarda-chuva",
     "R$ 12/conta/mês",
     AZUL,
     "Vale · B3 · LATAM · GOL\nGerdau · Braskem\ne outras 93 empresas"),
]

x = 0.4
for tier, modelo, qtd, criterio, tatica, budget, cor, exemplos in tiers:
    rect(slide, x, 1.05, 4.1, 5.55, AZUL_CLARO)
    rect(slide, x, 1.05, 4.1, 0.55, cor)
    txt(slide, tier,   x+0.12, 1.1,  2.0, 0.42, 14, bold=True, color=BRANCO if cor != AMARELO else AZUL)
    txt(slide, modelo, x+2.1,  1.12, 1.9, 0.38, 10, color=BRANCO if cor != AMARELO else AZUL)
    txt(slide, qtd,    x+0.12, 1.65, 3.8, 0.38, 20, bold=True, color=cor if cor != AMARELO else AZUL)
    div(slide, x+0.12, 2.1, 3.8, cor)
    txt(slide, "Critério", x+0.12, 2.2,  1.2, 0.3, 9, bold=True, color=AZUL)
    txt(slide, criterio,   x+0.12, 2.48, 3.8, 0.7, 10, color=CINZA_TEXTO)
    txt(slide, "Tática",   x+0.12, 3.22, 1.2, 0.3, 9, bold=True, color=AZUL)
    txt(slide, tatica,     x+0.12, 3.5,  3.8, 0.7, 10, color=CINZA_TEXTO)
    txt(slide, "Budget",   x+0.12, 4.25, 1.5, 0.3, 9, bold=True, color=AZUL)
    txt(slide, budget,     x+0.12, 4.5,  3.8, 0.35, 11, bold=True, color=cor if cor != AMARELO else AZUL)
    txt(slide, "Exemplos", x+0.12, 4.9, 1.5, 0.3, 9, bold=True, color=AZUL)
    txt(slide, exemplos,   x+0.12, 5.18, 3.8, 0.75, 9, italic=True, color=CINZA_TEXTO)
    x += 4.44

# ═══════════════════════════════════════
# SLIDE 4 — CONTAS TIER 1
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
header(slide, "Tier 1 — 7 Contas Estratégicas (Dores Identificadas no VoC)")

contas_t1 = [
    ("COMIGO",       "RFC + S4 pesado",     "RFC urgente — a fábrica pode parar",         VERMELHO),
    ("Geoavia",      "RFC em risco",        "RFC urgente — migração antes do bloqueio",    VERMELHO),
    ("Tuberfil",     "SAC + RFC + suporte", "Sem parceiro — quem sustenta o ambiente?",    VERDE),
    ("Grupo Soufer", "SAC parado",          "SAC parado — licença paga, nada em uso",      AMARELO),
    ("Docol",        "SAC parado",          "SAC parado — licença paga, nada em uso",      AMARELO),
    ("Algar",        "SAC parado",          "SAC parado + evolução para planning",         AZUL),
    ("BTG Pactual",  "S4 sobrecarregado",   "S4 pesado — cada query aparece na fatura",   LARANJA),
]

headers_row = ["Empresa", "Dor (VoC)", "Ângulo da mensagem", "Prioridade"]
col_x = [0.4, 3.0, 6.0, 11.0]
col_w = [2.5, 2.8, 4.8, 2.1]
rect(slide, 0.4, 1.05, 12.5, 0.45, AZUL)
for j, h in enumerate(headers_row):
    txt(slide, h, col_x[j]+0.1, 1.1, col_w[j], 0.35, 10, bold=True, color=AMARELO)

y = 1.52
for i, (empresa, dor, angulo, cor) in enumerate(contas_t1):
    bg = AZUL_CLARO if i % 2 == 0 else BRANCO
    rect(slide, 0.4, y, 12.5, 0.6, bg)
    rect(slide, 0.4, y, 0.06, 0.6, cor)
    txt(slide, empresa, col_x[0]+0.12, y+0.12, col_w[0], 0.36, 11, bold=True, color=AZUL)
    txt(slide, dor,     col_x[1]+0.1,  y+0.12, col_w[1], 0.36, 10, color=CINZA_TEXTO)
    txt(slide, angulo,  col_x[2]+0.1,  y+0.12, col_w[2], 0.36, 10, italic=True, color=CINZA_TEXTO)
    rect(slide, col_x[3]+0.1, y+0.1, 1.8, 0.38, cor)
    label = "RFC Urgente" if cor == VERMELHO else ("Sem Parceiro" if cor == VERDE else ("SAC Parado" if cor == AMARELO else ("S4 Pesado" if cor == LARANJA else "Suporte")))
    txt(slide, label, col_x[3]+0.1, y+0.1, 1.8, 0.38, 9, bold=True,
        color=AZUL if cor == AMARELO else BRANCO, align=PP_ALIGN.CENTER)
    y += 0.63

rect(slide, 0.4, y+0.1, 12.5, 0.5, AZUL_CLARO)
txt(slide, "Tática Tier 1: InMail personalizado por empresa + SDR outreach dia 3-5 após ativação + Sponsored Content segmentado. Budget: R$ 250/conta/mês.",
    0.55, y+0.15, 12.2, 0.38, 10, italic=True, color=AZUL)

# ═══════════════════════════════════════
# SLIDE 5 — MENSAGENS POR VERTICAL
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
header(slide, "Tier 2 — Mensagens por Vertical (50 Contas)")

verticais = [
    ("Agro /\nBioenergia", "25 contas",
     '"Na safra, seu S4 não pode travar. Cada relatório no transacional custa performance."',
     '"Desafogue seu S4 antes da safra"',
     "Atvos, São Martinho, Citrosuco, CTC, Coopercitrus, Copercana, Usina Coruripe, Tereos + 17 outras",
     VERDE),
    ("Energia /\nUtilities", "12 contas",
     '"A SAP vai bloquear extrações via RFC. Empresas de energia com arquiteturas legadas estão em alerta."',
     '"Sua extração via RFC está em risco"',
     "Cemig, CPFL, Eletrobras, Samarco, Petrobras, BRK Ambiental, Nexa Resources + 5 outras",
     AZUL),
    ("Indústria /\nManufatura", "13 contas",
     '"O time de consultoria foi embora depois da implantação. Quem evolui seu SAC agora?"',
     '"Sem parceiro técnico para o SAC?"',
     "Gerdau, Braskem, Embraer, Tupy, WEG, Tramontina, Klabin, Arauco, M. Dias Branco + 4 outras",
     LARANJA),
]

x = 0.4
for vertical, qtd, intro, headline, exemplos, cor in verticais:
    rect(slide, x, 1.05, 4.1, 5.55, AZUL_CLARO)
    rect(slide, x, 1.05, 4.1, 0.08, cor)
    txt(slide, vertical,  x+0.12, 1.15, 2.5, 0.55, 13, bold=True, color=AZUL)
    rect(slide, x+2.7, 1.15, 1.3, 0.42, cor)
    txt(slide, qtd, x+2.7, 1.15, 1.3, 0.42, 11, bold=True,
        color=BRANCO if cor != AMARELO else AZUL, align=PP_ALIGN.CENTER)
    div(slide, x+0.12, 1.72, 3.8, cor)
    txt(slide, "Intro text LinkedIn:", x+0.12, 1.82, 3.8, 0.28, 9, bold=True, color=AZUL)
    txt(slide, intro, x+0.12, 2.1, 3.8, 1.0, 10, italic=True, color=CINZA_TEXTO)
    txt(slide, "Headline:", x+0.12, 3.15, 3.8, 0.28, 9, bold=True, color=AZUL)
    rect(slide, x+0.12, 3.45, 3.8, 0.55, cor if cor != LARANJA else RGBColor(0xFF,0xF3,0xE0))
    txt(slide, headline, x+0.22, 3.5, 3.6, 0.45, 11, bold=True,
        color=AZUL if cor == LARANJA or cor == VERDE else BRANCO)
    txt(slide, "Empresas:", x+0.12, 4.1, 3.8, 0.28, 9, bold=True, color=AZUL)
    txt(slide, exemplos,   x+0.12, 4.38, 3.8, 1.1, 9, italic=True, color=CINZA_TEXTO)
    x += 4.44

# ═══════════════════════════════════════
# SLIDE 6 — ORQUESTRAÇÃO COM VENDAS
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
header(slide, "Orquestração Marketing + Vendas — ABM Tier 1")

txt(slide, "ABM funciona quando marketing e vendas estão sincronizados por conta.",
    0.5, 1.0, 12.3, 0.35, 11, italic=True, color=CINZA_TEXTO)

passos_orch = [
    ("Dia 0",    "Conta entra na audiência LinkedIn — campanha ativa",         "Marketing",  AZUL),
    ("Dia 3–5",  "SDR faz outreach personalizado com ângulo da dor (VoC)",     "Vendas",     VERDE),
    ("Sem 1–4",  "Conta vê anúncio 3–5x por semana (Sponsored + InMail)",      "Marketing",  AZUL),
    ("Contínuo", "Engajamento detectado (clique / visita) → alerta no HubSpot","Marketing",  AMARELO),
    ("< 24h",    "Follow-up de vendas imediato após sinal de engajamento",      "Vendas",     VERDE),
    ("30 dias",  "Revisão de contas frias — trocar ângulo ou pausar",          "Mkt + Vdas", CINZA_TEXTO),
]

y = 1.45
for timing, acao, resp, cor in passos_orch:
    rect(slide, 0.4, y, 12.5, 0.68, AZUL_CLARO if passos_orch.index((timing,acao,resp,cor)) % 2 == 0 else BRANCO)
    rect(slide, 0.4, y, 1.35, 0.68, cor if cor != CINZA_TEXTO else RGBColor(0x88,0x88,0x88))
    txt(slide, timing, 0.4, y+0.15, 1.35, 0.38, 11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(slide, acao,   1.85, y+0.15, 8.5,  0.38, 11, color=CINZA_TEXTO)
    rect(slide, 10.5, y+0.1, 2.25, 0.44, AZUL if resp == "Marketing" else (VERDE if resp == "Vendas" else AMARELO))
    txt(slide, resp, 10.5, y+0.1, 2.25, 0.44, 11, bold=True,
        color=BRANCO if resp != "Mkt + Vdas" else AZUL, align=PP_ALIGN.CENTER)
    y += 0.73

txt(slide, "Sinais de engajamento a monitorar: clique em anúncio ABM · visita à landing page · abertura de InMail · download de material",
    0.5, 6.45, 12.3, 0.38, 10, italic=True, color=CINZA_TEXTO)

# ═══════════════════════════════════════
# SLIDE 7 — METAS INTEGRADAS
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
header(slide, "Metas Q3 2026 — Campanha Integrada")

# Demand gen
txt(slide, "Geração de Demanda (inbound)", 0.5, 1.05, 6.0, 0.4, 13, bold=True, color=AZUL)
div(slide, 0.5, 1.48, 6.0, AMARELO)

metas_dg = [
    ("Leads gerados",        "8–12"),
    ("MQLs",                 "4–7"),
    ("Reuniões agendadas",   "2–4"),
    ("Oportunidades",        "1–2"),
    ("Pipeline inbound",     "R$ 678k–1,3M"),
]
y = 1.55
for i, (label, val) in enumerate(metas_dg):
    bg = CINZA_CLARO if i % 2 == 0 else BRANCO
    rect(slide, 0.5, y, 6.0, 0.48, bg)
    txt(slide, label, 0.65, y+0.1, 3.8, 0.28, 11, color=CINZA_TEXTO)
    txt(slide, val,   4.5, y+0.1, 1.9, 0.28, 11, bold=True, color=AZUL, align=PP_ALIGN.RIGHT)
    y += 0.5

# ABM
txt(slide, "ABM (por conta)", 7.2, 1.05, 5.7, 0.4, 13, bold=True, color=AZUL)
div(slide, 7.2, 1.48, 5.7, VERDE)

metas_abm = [
    ("Cobertura de conta",          "> 70% das 156"),
    ("Reuniões Tier 1",             "2–3 de 7 contas"),
    ("Contas Tier 2 engajadas",     "> 30%"),
    ("Oportunidades ABM",           "1–2"),
    ("Pipeline ABM",                "R$ 678k–2M"),
]
y = 1.55
for i, (label, val) in enumerate(metas_abm):
    bg = CINZA_CLARO if i % 2 == 0 else BRANCO
    rect(slide, 7.2, y, 5.7, 0.48, bg)
    txt(slide, label, 7.35, y+0.1, 3.5, 0.28, 11, color=CINZA_TEXTO)
    txt(slide, val,   10.9, y+0.1, 1.9, 0.28, 11, bold=True, color=VERDE, align=PP_ALIGN.RIGHT)
    y += 0.5

# Total consolidado
rect(slide, 0.5, 4.2, 12.4, 1.05, AZUL)
txt(slide, "Total consolidado Q3 2026", 0.7, 4.28, 5.0, 0.4, 13, bold=True, color=AMARELO)
totais = [("Oportunidades abertas", "2–4"), ("Pipeline gerado", "R$ 1,3M–3,3M")]
x = 6.5
for label, val in totais:
    txt(slide, label, x,   4.28, 3.0, 0.35, 11, color=BRANCO)
    txt(slide, val,   x,   4.65, 3.0, 0.45, 18, bold=True, color=AMARELO)
    x += 3.3

rect(slide, 0.5, 5.35, 12.4, 0.55, AZUL_CLARO)
txt(slide, "Racional: ticket médio R$ 678k — 1 deal fechado paga a campanha inteira. O objetivo do Q3 é criar o primeiro ciclo digital + ABM estruturado e escalar no Q4.",
    0.65, 5.42, 12.1, 0.42, 10, italic=True, color=AZUL)

# KPIs ABM
txt(slide, "KPIs ABM — medir por conta, não só por lead:", 0.5, 6.0, 12.4, 0.35, 10, bold=True, color=AZUL)
txt(slide, "Cobertura · Engajamento por empresa · Velocidade de pipeline · Win rate ABM vs geral",
    0.5, 6.35, 12.4, 0.35, 10, italic=True, color=CINZA_TEXTO)

# ═══════════════════════════════════════
# SLIDE 8 — BUDGET INTEGRADO
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
header(slide, "Budget Integrado — Q3 2026")

itens = [
    ("LinkedIn Ads — ABM Tier 1 (7 × R$250 × 3 meses)",      "R$ 5.250",  VERMELHO),
    ("LinkedIn Ads — ABM Tier 2 (50 × R$60 × 3 meses)",      "R$ 9.000",  AMARELO),
    ("LinkedIn Ads — ABM Tier 3 (99 × R$12 × 3 meses)",      "R$ 3.564",  AZUL),
    ("LinkedIn Ads — Demand gen (fora da lista ABM)",         "R$ 3.000–5.000", AZUL),
    ("Google Ads — RFC + SAC + Datasphere suporte",           "R$ 4.500–9.000", VERDE),
    ("Produção de peças (banners, LP, emails, InMails)",      "R$ 1.500–2.500", CINZA_TEXTO),
    ("TOTAL Q3 ESTIMADO",                                     "R$ 26.800–34.300", AZUL),
]

y = 1.1
for i, (item, valor, cor) in enumerate(itens):
    is_total = "TOTAL" in item
    bg = AZUL if is_total else (AZUL_CLARO if i % 2 == 0 else BRANCO)
    rect(slide, 0.5, y, 12.3, 0.58 if is_total else 0.5, bg)
    if not is_total:
        rect(slide, 0.5, y, 0.06, 0.5, cor if cor != CINZA_TEXTO else RGBColor(0x88,0x88,0x88))
    c_txt = BRANCO if is_total else CINZA_TEXTO
    c_val = AMARELO if is_total else AZUL
    txt(slide, item,  0.65, y+0.11, 9.0, 0.35, 11 if not is_total else 13, bold=is_total, color=c_txt)
    txt(slide, valor, 9.7,  y+0.11, 3.0, 0.35, 11 if not is_total else 13, bold=True, color=c_val if not is_total else AMARELO, align=PP_ALIGN.RIGHT)
    y += 0.54 if not is_total else 0.62

rect(slide, 0.5, y+0.1, 12.3, 0.75, AZUL_CLARO)
txt(slide, "Google Ads atual: CPC R$ 2,48 · CTR 4,97% — menor custo, maior previsibilidade. Escalar primeiro.",
    0.65, y+0.15, 12.0, 0.3, 10, italic=True, color=AZUL)
txt(slide, "ABM concentra budget onde já há relação. Menor risco, maior taxa de conversão esperada.",
    0.65, y+0.45, 12.0, 0.3, 10, italic=True, color=AZUL)

# ═══════════════════════════════════════
# SLIDE 9 — CRONOGRAMA
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
header(slide, "Cronograma Integrado — Q3 2026")

semanas = [
    ("Pré\naté 04/07", [
        "Exportar 156 contas → CSV → upload LinkedIn",
        "Criar audiências Matched (Tier 1, 2, 3)",
        "Criar InMails Tier 1 personalizados (7)",
        "Criar emails demand gen + landing page",
        "Briefar SDR/vendas — Tier 1 playbook",
    ]),
    ("Jul\nS1-S2\n07–18/07", [
        "Reativar 95 warm leads (WhatsApp + email RFC)",
        "Ativar LinkedIn Ads Tier 1 — SDR outreach",
        "Ativar Tier 2 Agro + Energia",
        "LinkedIn Ads demand gen + Google Ads",
        "Posts LinkedIn ângulo SAC parado",
    ]),
    ("Jul\nS3-S4\n21–31/07", [
        "Email RFC sequência — demand gen",
        "Ativar Tier 2 Indústria + Tier 3",
        "Artigo blog: RFC bloqueado",
        "Revisar cobertura de conta ABM",
        "Posts ângulo sem parceiro técnico",
    ]),
    ("Ago\nS5-S8\n04–29/08", [
        "Email S4 sobrecarregado + SAC como BI",
        "Revisão Tier 1: quem engajou?",
        "Trocar variante menos performática",
        "Artigo: quem sustenta seu SAC",
        "Revisão mid-campaign completa (25/08)",
    ]),
    ("Set\nS9-S12\n01–30/09", [
        "Acelerar canal líder (definir em 25/08)",
        "Reativar leads não convertidos",
        "Qualificação comercial — contas engajadas",
        "Relatório Q3 + aprendizados",
        "Planejamento Q4 com base nos dados",
    ]),
]

x = 0.35
for periodo, ativs in semanas:
    rect(slide, x, 1.05, 2.5, 0.6, AZUL)
    txt(slide, periodo, x, 1.05, 2.5, 0.6, 10, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    rect(slide, x, 1.65, 2.5, 5.0, AZUL_CLARO)
    y = 1.72
    for at in ativs:
        txt(slide, f"• {at}", x+0.1, y, 2.3, 0.55, 9, color=CINZA_TEXTO)
        y += 0.57
    x += 2.6

# ═══════════════════════════════════════
# SLIDE 10 — PRÓXIMOS PASSOS
# ═══════════════════════════════════════
slide = prs.slides.add_slide(blank())
rect(slide, 0, 0, 13.33, 7.5, AZUL)
rect(slide, 0, 0, 13.33, 0.12, AMARELO)
rect(slide, 0, 6.88, 13.33, 0.62, AMARELO)

txt(slide, "Próximos Passos — Esta Semana", 0.8, 0.28, 11, 0.6, 26, bold=True, color=BRANCO)
div(slide, 0.8, 0.98, 11.7, AMARELO)

passos = [
    ("1", "Reativar 95 warm leads via WhatsApp — ângulo RFC urgente — custo zero",        "Agora", VERMELHO),
    ("2", "Exportar lista 156 empresas do Leads_Geral.xlsx em CSV para LinkedIn",          "Até 01/07", AMARELO),
    ("3", "Upload no LinkedIn Campaign Manager — criar 4 audiências Matched por tier",     "Até 01/07", AMARELO),
    ("4", "Criar InMails personalizados para 7 contas Tier 1 (dor por empresa)",           "Até 04/07", AMARELO),
    ("5", "Briefar SDR/vendas: conta, dor identificada no VoC, ângulo, timing",           "Até 04/07", AMARELO),
    ("6", "Criar emails demand gen + landing page guarda-chuva",                           "Até 04/07", AMARELO),
    ("7", "Ativar LinkedIn Ads Tier 1 + Google Ads na semana 1 de julho",                 "07/07", VERDE),
    ("8", "Primeira revisão de cobertura ABM e engajamento Tier 1",                       "25/07", VERDE),
]

y = 1.1
for num, acao, prazo, cor in passos:
    rect(slide, 0.8, y, 0.55, 0.52, cor)
    txt(slide, num, 0.8, y, 0.55, 0.52, 14, bold=True,
        color=AZUL if cor == AMARELO else BRANCO, align=PP_ALIGN.CENTER)
    txt(slide, acao,  1.45, y+0.1, 9.8, 0.32, 11, color=BRANCO)
    txt(slide, prazo, 11.35, y+0.1, 1.6, 0.32, 11, bold=True, color=AMARELO, align=PP_ALIGN.RIGHT)
    y += 0.6

txt(slide, "solveplan", 0.8, 6.96, 4, 0.38, 13, bold=True, color=AZUL)

# ═══════════════════════════════════════
# SALVAR
# ═══════════════════════════════════════
output = (
    r"c:\Users\franc\solveplan.com\Roberto Molina - Marketing"
    r"\1. MKT Estrategy\3. Agentes de IA\ccos-ratos"
    r"\marketing\campanhas\fabrica-de-analytics-q3-2026"
    r"\Fabrica_Analytics_ABM_Q3_2026.pptx"
)
prs.save(output)
print(f"PPT salvo: {output}")
